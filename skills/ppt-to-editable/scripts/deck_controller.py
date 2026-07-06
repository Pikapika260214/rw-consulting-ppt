#!/usr/bin/env python
"""Deck controller for ppt-to-editable v3 two-mode preview.

This script owns the deterministic deck-level work:

- validate a PPTX as image-only;
- extract unchanged per-slide source images;
- write deck and slide job manifests;
- create per-slide Agent task briefs;
- optionally assemble unresolved slides as image-fallback pages with a sticker.

Per-slide workers must use this package's own scripts/ and references/
directories as the primary single-page engine.
"""

from __future__ import annotations

import argparse
import datetime as dt
import hashlib
import json
import os
import posixpath
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path
import xml.etree.ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from visual_strategy_rules import (  # noqa: E402
    full_slide_textless_background_issue,
    high_risk_native_strategy_issues,
    region_crop_plan_issues,
    zero_crop_high_native_issue,
)


NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

DEFAULT_STICKER_TEXT = "可编辑转换失败"
OCR_TEXT_USABLE_STATUS = "passed-text-usable"
OCR_GEOMETRY_ONLY_STATUS = "passed-geometry-only"
OCR_FAILED_STATUS = "failed"

GATES_SCHEMA_VERSION = "ppt-to-editable-gates-v3"
PROBE_SCHEMA_VERSION = "ppt-to-editable-deck-probe-v1"
CONVERSION_MODE_CHOICES = (
    "single-page-token-saving",
    "multi-agent-high-quality",
    "full-deck-convenience",
)
RETIRED_CONVERSION_MODE_CHOICES = ("multi-page-sequential-token-saving",)
SCOPE_GATE_CHOICES = ("all", "selected", "sample")
WORKER_MODE_CHOICES = ("app-native", "external-codex", "manual")
RETIRED_WORKER_MODE_CHOICES = ("sequential",)

REQUIRED_PROGRESS_STATUSES = [
    "started",
    "ocr_done",
    "plan_done",
    "packaging_done",
    "qa_done",
]

BLOCKING_SUCCESS_LIMITATION_PATTERNS = [
    "ocr-unavailable",
    "ocr unavailable",
    "manual-layout-guess",
    "manual layout guess",
    "manual bbox",
    "manual text review",
    "text smaller than source",
    "body text smaller than source",
    "small central bullet text compared with source",
    "evidence card text is smaller than source",
    "layout drift",
    "major layout drift",
    "native-redraw-worsened-fidelity",
    "rendered-text-overflow",
    "rendered-text-overlap",
    "font-fallback",
    "font fallback",
    "unwanted-shape-shadow",
    "source-crop-text-residue",
    "tiny-font-outlier",
    "font-size-hierarchy-drift",
    "unnecessary-font-shrink",
    "invented-text-background",
    "pale rectangle",
    "text backing rectangle",
    "backing rectangles create",
    "native-redraw-complex-artwork",
    "complex artwork native redraw drift",
    "missing-non-text-anchor",
    "missing-arrow",
    "missing-arrowhead",
    "erased-arrow",
    "broken-route-line",
    "broken-connector",
    "missing-connector",
    "missing-step-marker",
    "missing-card-edge",
    "cropped-icon",
    "icon stroke missing",
    "full-slide textless",
    "full slide textless",
    "whole-slide textless",
    "whole slide textless",
    "textless full-slide",
    "textless full slide",
    "full-slide visual layer",
    "full slide visual layer",
    "whole-slide image",
    "full-slide background disguised",
    "full_slide_textless",
]

PRIMARY_SINGLE_PAGE_ENGINE_ID = "ppt-to-editable-v3-1-preview"
RECONSTRUCTION_PACKAGING_SCRIPTS = {
    "run_reconstruction_qa.py",
    "package_reconstruction_deck.py",
}
CLEAN_BACKGROUND_PACKAGING_SCRIPTS = {
    "package_clean_background_deck.py",
}


def single_page_engine_root() -> Path:
    return Path(__file__).resolve().parents[1]


def single_page_engine_script_paths() -> dict[str, str]:
    scripts_dir = single_page_engine_root() / "scripts"
    return {
        "preflight_ocr_runtime": str(scripts_dir / "preflight_ocr_runtime.py"),
        "run_rapidocr": str(scripts_dir / "run_rapidocr.py"),
        "ocr_to_reconstruction_plan_scaffold": str(scripts_dir / "ocr_to_reconstruction_plan_scaffold.py"),
        "extract_source_crops": str(scripts_dir / "extract_source_crops.py"),
        "make_textless_crops": str(scripts_dir / "make_textless_crops.py"),
        "make_clean_badge_composites": str(scripts_dir / "make_clean_badge_composites.py"),
        "run_reconstruction_qa": str(scripts_dir / "run_reconstruction_qa.py"),
        "package_reconstruction_deck": str(scripts_dir / "package_reconstruction_deck.py"),
        "package_clean_background_deck": str(scripts_dir / "package_clean_background_deck.py"),
        "inspect_pptx_editability": str(scripts_dir / "inspect_pptx_editability.py"),
    }


def context_reference_paths() -> dict[str, str]:
    engine_root = single_page_engine_root()
    refs_dir = engine_root / "references"
    return {
        "worker_compact_protocol_template": str(refs_dir / "worker-compact-protocol.md"),
        "single_page_engine_skill_3_1": str(engine_root / "SKILL.md"),
        "single_page_engine_source_crop_tools": str(refs_dir / "source-crop-tools.md"),
        "single_page_engine_reconstruction_qa": str(refs_dir / "reconstruction-qa.md"),
        "single_page_engine_reconstruction_plan_schema": str(refs_dir / "source-reconstruction-plan-schema.md"),
        "single_page_engine_ocr_to_plan_review": str(refs_dir / "ocr-to-plan-review-guide.md"),
        "active_skill_fallback": str(Path.home() / ".codex" / "skills" / "ppt-to-editable" / "SKILL.md"),
        "visual_fidelity_contract": str(refs_dir / "visual-fidelity-contract.md"),
        "visual_fidelity_execution": str(refs_dir / "visual-fidelity-execution.md"),
        "multi_agent_slide_protocol": str(refs_dir / "multi-agent-slide-protocol.md"),
    }


def make_page_conversion_contract(slide_index: int, source_size: dict) -> dict:
    """Small per-slide execution lock for worker decisions.

    This deliberately repeats only the high-impact defaults that should not be
    re-decided from long prose on every slide.
    """
    return {
        "schema_version": "ppt-to-editable-page-conversion-contract-v1",
        "slide_index": slide_index,
        "source_size": source_size,
        "default_font_family": "Microsoft YaHei",
        "default_shape_effects": "none",
        "default_text_box_fill": "transparent_unless_source_has_visible_text_background",
        "fallback_decision": "prefer_fidelity_over_extra_native_shapes",
        "text_rules": {
            "meaningful_text_native_when_practical": True,
            "review_ocr_before_final_text": True,
            "preserve_multiline_semantic_blocks": True,
            "conservative_multiline_grouping": {
                "merge_when_same_font_size_weight_color_alignment_region": True,
                "do_not_merge_mixed_style_labels_tables_legends_axes_badges_steps_or_cross_column_text": True,
                "mark_intentional_line_level_with_line_level_trace_required_or_do_not_merge": True,
            },
            "avoid_single_character_wraps": True,
            "do_not_shrink_body_or_callout_text_into_tiny_outliers": True,
        },
        "shape_rules": {
            "rebuild_simple_flat_geometry": True,
            "clear_unwanted_effects": True,
            "no_default_shadow_glow_bevel_reflection_or_soft_edge": True,
            "do_not_add_invented_text_backing_fills": True,
        },
        "visual_rules": {
            "complex_visuals_crop_first": True,
            "source_specific_modules_use_crop_when_native_redraw_would_drift": True,
            "preserve_non_text_anchors_during_text_removal": True,
            "whole_slide_textless_crop_requires_explicit_fallback_acceptance": True,
            "single_full_slide_picture_cannot_satisfy_region_crop_plan": True,
        },
        "success_reporting_fields": [
            "font_policy_applied",
            "effects_cleared",
            "font_size_consistency_review",
            "complex_visual_strategy_review",
            "text_background_policy",
            "non_text_anchor_preservation_review",
        ],
    }


def make_crop_manifest_contract(slide_index: int, source_size: dict) -> dict:
    """Contract for crop planning before packaging."""
    return {
        "schema_version": "ppt-to-editable-crop-manifest-contract-v1",
        "slide_index": slide_index,
        "source_size": source_size,
        "artifact": "crop_manifest.json",
        "purpose": (
            "Make crop decisions explicit before packaging so complex visuals are not "
            "accidentally redrawn or erased during text removal."
        ),
        "allowed_region_strategies": [
            "native_text",
            "native_shape",
            "source_crop",
            "textless_crop_plus_editable_text",
            "clean_badge_composite",
            "image_fallback",
        ],
        "semantic_region_guidance": {
            "complex_slide_min_regions": 3,
            "typical_complex_slide_regions": "3-5",
            "whole_slide_crop_policy": "fallback_only_with_explicit_user_acceptance",
        },
        "review_gates": [
            "ocr_text_overlap_reviewed",
            "source_text_residue_absent_or_explained",
            "crop_boxes_complete_and_centered",
            "must_not_erase_non_text_anchors",
            "text_overlay_present_for_textless_text_regions",
        ],
        "tool_manifest_aliases": {
            "source_crop_manifest": "source-crops.json",
            "textless_crop_manifest": "textless-crops.json",
            "clean_badge_composite_manifest": "clean-badge-composites.json",
        },
    }


def routing_decision_table() -> list[dict]:
    return [
        {
            "id": "pure_text_to_native_text",
            "when": "Region is meaningful visible text without complex visual background.",
            "route": "native_text",
            "quality_bias": "editable",
        },
        {
            "id": "simple_flat_geometry_to_native_shape",
            "when": "Region is a flat rectangle, line, divider, circle, simple badge, simple arrow, or plain table grid.",
            "route": "native_shape",
            "quality_bias": "editable_when_faithful",
        },
        {
            "id": "complex_visual_to_source_crop",
            "when": "Region is a product image, detailed icon, shaded badge, gradient bar, curved path, textured module, or source-specific artwork.",
            "route": "source_crop",
            "quality_bias": "fidelity",
        },
        {
            "id": "text_in_complex_visual_to_textless_crop_overlay",
            "when": "Complex visual region contains readable text that should remain editable.",
            "route": "textless_crop_plus_editable_text",
            "quality_bias": "fidelity_plus_editable_text",
        },
        {
            "id": "ocr_unusable_to_fallback",
            "when": "OCR is unavailable, geometry-only, or text cannot be reviewed reliably.",
            "route": "failed_retryable_or_image_fallback",
            "quality_bias": "do_not_fake_quality_conversion",
        },
    ]


def ensure_worker_compact_protocol(run_dir: Path) -> Path:
    template_path = single_page_engine_root() / "references" / "worker-compact-protocol.md"
    output_path = run_dir / "WORKER_COMPACT_PROTOCOL.md"
    if not template_path.exists():
        raise FileNotFoundError(f"Compact worker protocol template missing: {template_path}")
    text = template_path.read_text(encoding="utf-8")
    compiled = (
        "<!-- Compiled for this ppt-to-editable run. Do not edit per slide. -->\n\n"
        + text
    )
    output_path.write_text(compiled, encoding="utf-8")
    return output_path


def fallback_reading_entries(reference_paths: dict[str, str]) -> list[dict]:
    return [
        {
            "id": "single_page_engine_skill_3_1",
            "path": reference_paths["single_page_engine_skill_3_1"],
            "read_when_blocked_or_needed": True,
            "reason": "Full skill body for mode, gate, or protocol ambiguity.",
        },
        {
            "id": "single_page_engine_source_crop_tools",
            "path": reference_paths["single_page_engine_source_crop_tools"],
            "read_when_blocked_or_needed": True,
            "reason": "Long source crop, textless crop, and clean badge composite rules.",
        },
        {
            "id": "single_page_engine_reconstruction_qa",
            "path": reference_paths["single_page_engine_reconstruction_qa"],
            "read_when_blocked_or_needed": True,
            "reason": "Long reconstruction QA reference for failed or unclear QA gates.",
        },
        {
            "id": "single_page_engine_reconstruction_plan_schema",
            "path": reference_paths["single_page_engine_reconstruction_plan_schema"],
            "read_when_blocked_or_needed": True,
            "reason": "Full source_reconstruction_plan.json schema.",
        },
        {
            "id": "single_page_engine_ocr_to_plan_review",
            "path": reference_paths["single_page_engine_ocr_to_plan_review"],
            "read_when_blocked_or_needed": True,
            "reason": "OCR scaffold review rules.",
        },
        {
            "id": "visual_fidelity_contract",
            "path": reference_paths["visual_fidelity_contract"],
            "read_when_blocked_or_needed": True,
            "reason": "Long visual fidelity preference contract.",
        },
        {
            "id": "visual_fidelity_execution",
            "path": reference_paths["visual_fidelity_execution"],
            "read_when_blocked_or_needed": True,
            "reason": "Long execution-time fidelity reference.",
        },
        {
            "id": "multi_agent_slide_protocol",
            "path": reference_paths["multi_agent_slide_protocol"],
            "read_when_blocked_or_needed": True,
            "reason": "Long multi-agent handoff and worker identity protocol.",
        },
    ]


def reference_loading_policy(required_reading: list[dict], fallback_reading: list[dict] | None = None) -> dict:
    return {
        "default_mode": "compact-protocol-first",
        "must_read_ids": [item["id"] for item in required_reading if item.get("must_read")],
        "fallback_reading_ids": [
            item["id"]
            for item in (fallback_reading or [])
            if item.get("read_when_blocked_or_needed")
        ],
        "conditional_reading_ids": [
            "ocr_troubleshooting",
            "clean_background_hybrid_contract",
            "text_layout_manifest",
            "failure_recovery_notes",
        ],
        "forbidden_context": [
            "old_round_outputs",
            "historical_test_reports",
            "prior_visual_qa_html_from_other_runs",
            "old_baseline_materials_unless_user_explicitly_requests_reuse",
            "external_development_folders_as_primary_engine",
        ],
        "principle": (
            "Read the compact protocol, slide job, and the two small per-slide "
            "contracts first. Read long references only when blocked, when a QA "
            "gate fails, or when the compact protocol is ambiguous. Do not load old "
            "rounds, process audits, or previous test artifacts as conversion "
            "context for a fresh run."
        ),
    }


def slide_sort_key(name: str) -> int:
    match = re.search(r"slide(\d+)\.xml$", name)
    return int(match.group(1)) if match else 0


def parse_slide_selection(selection: str | None) -> list[int] | None:
    if not selection:
        return None

    slides: list[int] = []
    for raw_part in selection.split(","):
        part = raw_part.strip()
        if not part:
            continue
        if "-" in part:
            start_text, end_text = [item.strip() for item in part.split("-", 1)]
            if not start_text.isdigit() or not end_text.isdigit():
                raise ValueError(f"Invalid slide range: {part}")
            start = int(start_text)
            end = int(end_text)
            if start <= 0 or end <= 0 or start > end:
                raise ValueError(f"Invalid slide range: {part}")
            slides.extend(range(start, end + 1))
        else:
            if not part.isdigit():
                raise ValueError(f"Invalid slide number: {part}")
            slide = int(part)
            if slide <= 0:
                raise ValueError(f"Invalid slide number: {part}")
            slides.append(slide)

    deduped: list[int] = []
    seen: set[int] = set()
    for slide in slides:
        if slide not in seen:
            deduped.append(slide)
            seen.add(slide)
    if not deduped:
        raise ValueError("Slide selection is empty")
    return deduped


def filter_slides_by_selection(slides: list[dict], selected: list[int] | None) -> list[dict]:
    if selected is None:
        return slides

    by_index = {slide["slide_index"]: slide for slide in slides}
    missing = [slide for slide in selected if slide not in by_index]
    if missing:
        available = [slide["slide_index"] for slide in slides]
        raise ValueError(f"Selected slides do not exist: {missing}; available slides: {available}")
    return [by_index[slide] for slide in selected]


def require_explicit_scope(slides: str | None, all_slides: bool) -> list[int] | None:
    if slides and all_slides:
        raise ValueError("Scope + Worker Gate violation: pass either --all-slides or --slides, not both.")
    if not slides and not all_slides:
        raise ValueError(
            "Scope + Worker Gate required: confirm the conversion scope with the user, then pass "
            "--all-slides or --slides \"1,3,5\". Do not default a multi-page PPTX to all pages."
        )
    return parse_slide_selection(slides)


def require_ocr_text_usable(ocr_runtime_status: str, allow_ocr_not_ready: bool) -> None:
    if ocr_runtime_status == OCR_TEXT_USABLE_STATUS:
        return
    if allow_ocr_not_ready:
        return
    raise RuntimeError(
        "OCR Runtime Gate required: OCR is not passed-text-usable. Stop and ask the user "
        "whether to set up OCR Runtime before conversion. Do not create worker jobs, reuse "
        "old baselines, or deliver a key-text-only fallback as a quality conversion."
    )


def gates_file_template() -> dict:
    return {
        "schema_version": GATES_SCHEMA_VERSION,
        "source_pptx_sha256": "<copy source_pptx_sha256 from the --probe output>",
        "conversion_mode_gate": {
            "user_choice": "single-page-token-saving | multi-agent-high-quality",
            "user_response_quote": "",
        },
        "ocr_runtime_gate": {
            "probe_ocr_runtime_status": "passed-text-usable | passed-geometry-only | failed",
            "setup_required": False,
            "auto_passed_existing_runtime": False,
            "explained_to_user": False,
            "user_confirmed_setup": False,
            "user_accepted_limited_run": False,
            "user_response_quote": "",
        },
        "scope_and_worker_gate": {
            "user_choice": "all | selected | sample",
            "slides": "1,3,5 (required when user_choice is selected or sample)",
            "worker_mode": "app-native | external-codex | manual",
            "worker_execution_approved": False,
            "external_codex_worker_approved": False,
            "token_cost_acknowledged": False,
            "content_sharing_acknowledged": False,
            "user_response_quote": "",
        },
    }


def load_and_validate_gates(
    gates_file: Path | None,
    source_pptx: Path,
    slides: str | None,
    all_slides: bool,
    allow_ocr_not_ready: bool,
) -> tuple[dict, Path]:
    """Validate the user-confirmation record before any run initialization.

    The gates file is the auditable evidence that conversion mode, OCR setup,
    slide scope, and per-slide worker/token authorization were actually handled.
    Initialization must fail without it, and must fail when it contradicts the
    CLI flags.
    """
    if gates_file is None:
        raise ValueError(
            "Gates file required: run --probe first, handle the Conversion Mode Gate, "
            "OCR Runtime Gate, and combined Scope + Worker Gate, record the answers in a gates file, then "
            "initialize with --gates-file. Do not initialize a run before all gates are "
            "confirmed or auto-recorded."
        )
    gates_file = gates_file.resolve()
    if not gates_file.exists():
        raise FileNotFoundError(f"Gates file not found: {gates_file}")
    gates = read_json(gates_file)
    if gates.get("schema_version") != GATES_SCHEMA_VERSION:
        raise ValueError(
            f"Gates file schema_version must be {GATES_SCHEMA_VERSION}; "
            f"got {gates.get('schema_version')!r}."
        )
    recorded_sha = gates.get("source_pptx_sha256")
    actual_sha = sha256_file(source_pptx)
    if recorded_sha != actual_sha:
        raise ValueError(
            "Gates file does not match this deck: source_pptx_sha256 differs from the input "
            "PPTX. Run --probe on this PPTX, handle the Conversion Mode Gate, OCR Runtime Gate, and combined "
            "Scope + Worker Gate again, and write a gates file for this deck."
        )

    conversion_gate = gates.get("conversion_mode_gate")
    if not isinstance(conversion_gate, dict):
        raise ValueError(
            "Conversion Mode Gate violation: record conversion_mode_gate with the user's "
            "choice between single-page-token-saving and multi-agent-high-quality."
        )
    conversion_mode = conversion_gate.get("user_choice")
    if conversion_mode in RETIRED_CONVERSION_MODE_CHOICES:
        raise ValueError(
            "Conversion Mode Gate violation: multi-page-sequential-token-saving is retired "
            "and not supported in this release. Use single-page-token-saving for one PNG/sample "
            "slide, or multi-agent-high-quality for multi-page PPTX conversion."
        )
    if conversion_mode not in CONVERSION_MODE_CHOICES:
        raise ValueError(
            "Conversion Mode Gate violation: conversion_mode_gate.user_choice must be one "
            f"of {CONVERSION_MODE_CHOICES}; got {conversion_mode!r}."
        )
    if not str(conversion_gate.get("user_response_quote") or "").strip():
        raise ValueError(
            "Conversion Mode Gate violation: record the user's actual answer in "
            "conversion_mode_gate.user_response_quote."
        )
    if conversion_mode == "single-page-token-saving":
        selected_for_mode = None if all_slides else parse_slide_selection(slides)
        if all_slides or selected_for_mode is None or len(selected_for_mode) != 1:
            raise ValueError(
                "Conversion Mode Gate mismatch: single-page-token-saving mode must initialize "
                "exactly one selected slide with --slides, not --all-slides."
            )

    ocr_gate = gates.get("ocr_runtime_gate") or {}
    if not isinstance(ocr_gate, dict):
        raise ValueError("OCR Runtime Gate violation: ocr_runtime_gate must be an object.")
    setup_required = ocr_gate.get("setup_required") is True
    auto_passed_existing_runtime = ocr_gate.get("auto_passed_existing_runtime") is True
    probe_ocr_runtime_status = normalize_ocr_runtime_status(
        ocr_gate.get("probe_ocr_runtime_status"),
        None,
    )
    if setup_required:
        if ocr_gate.get("explained_to_user") is not True:
            raise ValueError(
                "OCR Runtime Gate violation: ocr_runtime_gate.explained_to_user must be true "
                "when OCR setup is required. Explain what OCR provides and why first-time "
                "setup may take time before initializing a run."
            )
        if not str(ocr_gate.get("user_response_quote") or "").strip():
            raise ValueError(
                "OCR Runtime Gate violation: record the user's actual answer in "
                "ocr_runtime_gate.user_response_quote when setup is required."
            )
        if (
            ocr_gate.get("user_confirmed_setup") is not True
            and ocr_gate.get("user_accepted_limited_run") is not True
        ):
            raise ValueError(
                "OCR Runtime Gate violation: setup_required=true requires either "
                "ocr_runtime_gate.user_confirmed_setup=true or "
                "ocr_runtime_gate.user_accepted_limited_run=true."
            )
    elif not auto_passed_existing_runtime and probe_ocr_runtime_status != OCR_TEXT_USABLE_STATUS:
        raise ValueError(
            "OCR Runtime Gate violation: setup_required=false requires either "
            "ocr_runtime_gate.auto_passed_existing_runtime=true or "
            "ocr_runtime_gate.probe_ocr_runtime_status=passed-text-usable."
        )
    if allow_ocr_not_ready and ocr_gate.get("user_accepted_limited_run") is not True:
        raise ValueError(
            "--user-accepted-ocr-limited-run requires ocr_runtime_gate.user_accepted_limited_run "
            "to be true in the gates file, recorded from the user's explicit acceptance of a "
            "diagnostic or fallback-only run."
        )

    scope_worker_gate = gates.get("scope_and_worker_gate")
    if not isinstance(scope_worker_gate, dict):
        raise ValueError(
            "Scope + Worker Gate violation: record scope_and_worker_gate with the user's "
            "page-scope choice, worker approval, token-cost acknowledgement, and content "
            "sharing acknowledgement."
        )
    choice = scope_worker_gate.get("user_choice")
    if choice not in SCOPE_GATE_CHOICES:
        raise ValueError(
            "Scope + Worker Gate violation: scope_and_worker_gate.user_choice must be one of "
            f"{SCOPE_GATE_CHOICES}; got {choice!r}."
        )
    worker_mode = scope_worker_gate.get("worker_mode")
    if worker_mode in RETIRED_WORKER_MODE_CHOICES:
        raise ValueError(
            "Scope + Worker Gate violation: worker_mode=sequential is retired and not supported "
            "in this release. Use app-native/manual per-slide workers for multi-agent-high-quality."
        )
    if worker_mode not in WORKER_MODE_CHOICES:
        raise ValueError(
            "Scope + Worker Gate violation: scope_and_worker_gate.worker_mode must be one "
            f"of {WORKER_MODE_CHOICES}; got {worker_mode!r}."
        )
    if scope_worker_gate.get("worker_execution_approved") is not True:
        raise ValueError(
            "Scope + Worker Gate violation: scope_and_worker_gate.worker_execution_approved "
            "must be true before creating per-slide worker jobs."
        )
    if scope_worker_gate.get("token_cost_acknowledged") is not True:
        raise ValueError(
            "Scope + Worker Gate violation: scope_and_worker_gate.token_cost_acknowledged "
            "must be true because more selected pages means more worker/token cost."
        )
    if scope_worker_gate.get("content_sharing_acknowledged") is not True:
        raise ValueError(
            "Scope + Worker Gate violation: scope_and_worker_gate.content_sharing_acknowledged "
            "must be true before per-page task content can be sent to worker execution."
        )
    if (
        worker_mode == "external-codex"
        and scope_worker_gate.get("external_codex_worker_approved") is not True
    ):
        raise ValueError(
            "Scope + Worker Gate violation: worker_mode=external-codex requires "
            "scope_and_worker_gate.external_codex_worker_approved=true."
        )
    if not str(scope_worker_gate.get("user_response_quote") or "").strip():
        raise ValueError(
            "Scope + Worker Gate violation: record the user's actual answer in "
            "scope_and_worker_gate.user_response_quote."
        )
    if choice == "all":
        if not all_slides:
            raise ValueError(
                "Scope + Worker Gate mismatch: gates file records user_choice=all; initialize with "
                "--all-slides."
            )
        if conversion_mode == "single-page-token-saving":
            raise ValueError(
                "Conversion Mode Gate mismatch: single-page-token-saving mode cannot use "
                "scope_and_worker_gate.user_choice=all."
            )
    else:
        if all_slides or not slides:
            raise ValueError(
                f"Scope + Worker Gate mismatch: gates file records user_choice={choice}; initialize "
                "with --slides matching scope_and_worker_gate.slides."
            )
        recorded_slides = scope_worker_gate.get("slides")
        if isinstance(recorded_slides, list):
            recorded = sorted({int(v) for v in recorded_slides})
        else:
            recorded = parse_slide_selection(str(recorded_slides or "") or None)
        requested = parse_slide_selection(slides)
        if conversion_mode == "single-page-token-saving" and len(requested) != 1:
            raise ValueError(
                "Conversion Mode Gate mismatch: single-page-token-saving mode must select exactly one slide."
            )
        if conversion_mode in {
            "full-deck-convenience",
            "multi-agent-high-quality",
        } and choice == "sample":
            raise ValueError(
                "Conversion Mode Gate mismatch: sample scope should use single-page-token-saving mode."
            )
        if recorded != requested:
            raise ValueError(
                f"Scope + Worker Gate mismatch: --slides resolves to {requested} but the gates file "
                f"records {recorded}. Re-confirm the scope with the user."
            )
        if choice == "sample" and len(requested) != 1:
            raise ValueError(
                "Scope + Worker Gate violation: user_choice=sample must select exactly one slide."
            )
    return gates, gates_file


def probe_deck(
    source_pptx: Path,
    ocr_python: Path | None = None,
    ocr_deps: Path | None = None,
) -> dict:
    """Read-only gate probe. Creates no run directory, extracts no files,
    installs and downloads nothing.

    Its output gives the agent everything needed to handle OCR setup only when
    needed and ask the plain-language page-scope question with token/content notes.
    """
    source_pptx = source_pptx.resolve()
    if not source_pptx.exists():
        raise FileNotFoundError(source_pptx)
    if source_pptx.suffix.lower() != ".pptx":
        raise ValueError("Only .pptx input is supported by the deck controller")

    with zipfile.ZipFile(source_pptx) as zf:
        slide_width_emu, slide_height_emu = read_slide_size(zf)
        slide_names = sorted(
            [name for name in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)],
            key=slide_sort_key,
        )
        if not slide_names:
            raise ValueError("No slides found in PPTX")
        part_names = set(zf.namelist())
        slide_probes = []
        for idx, slide_name in enumerate(slide_names, start=1):
            slide_xml = ET.fromstring(zf.read(slide_name))
            rels_name = f"ppt/slides/_rels/slide{idx}.xml.rels"
            rel_map = {}
            if rels_name in part_names:
                rels_xml = ET.fromstring(zf.read(rels_name))
                for rel in rels_xml.findall("rel:Relationship", NS):
                    rid = rel.attrib.get("Id")
                    target = rel.attrib.get("Target")
                    if rid and target:
                        rel_map[rid] = resolve_rel_target(slide_name, target)
            pics = slide_xml.findall(".//p:pic", NS)
            full_slide_pics = 0
            for pic in pics:
                blip = pic.find(".//a:blip", NS)
                rid = blip.attrib.get(f"{{{NS['r']}}}embed") if blip is not None else None
                xfrm = pic.find(".//a:xfrm", NS)
                off = xfrm.find("a:off", NS) if xfrm is not None else None
                ext = xfrm.find("a:ext", NS) if xfrm is not None else None
                target = rel_map.get(rid) if rid else None
                if (
                    off is not None
                    and ext is not None
                    and int(off.attrib.get("x", 0)) == 0
                    and int(off.attrib.get("y", 0)) == 0
                    and int(ext.attrib.get("cx", 0)) == slide_width_emu
                    and int(ext.attrib.get("cy", 0)) == slide_height_emu
                    and target in part_names
                ):
                    full_slide_pics += 1
            slide_probes.append(
                {
                    "slide_index": idx,
                    "picture_count": len(pics),
                    "is_single_full_slide_picture": len(pics) == 1 and full_slide_pics == 1,
                }
            )

    all_image_only = all(s["is_single_full_slide_picture"] for s in slide_probes)
    ocr_status = get_ocr_dependency_status(ocr_python, ocr_deps)
    ocr_runtime_status = normalize_ocr_runtime_status(
        ocr_status.get("ocr_runtime_status"),
        ocr_status.get("status"),
    )
    ocr_ready = ocr_runtime_status == OCR_TEXT_USABLE_STATUS
    slide_count = len(slide_probes)
    setup_script = Path(__file__).with_name("setup_ocr_runtime.py")
    return {
        "schema_version": PROBE_SCHEMA_VERSION,
        "source_pptx": str(source_pptx),
        "source_pptx_sha256": sha256_file(source_pptx),
        "slide_count": slide_count,
        "slide_indices": [s["slide_index"] for s in slide_probes],
        "all_slides_are_full_page_images": all_image_only,
        "slides": slide_probes,
        "ocr_runtime_status": ocr_runtime_status,
        "ocr_runtime_report": ocr_status,
        "user_confirmation_required": {
            "conversion_mode_gate": {
                "question": (
                    "你可以选择两种转换方式：\n\n"
                    "1. 单页省 token 模式\n"
                    "适合：你只想先转一页、测试效果，或者想继续发单张 PNG 变成可编辑 PPT。\n"
                    "消耗：较低。\n"
                    "结果：输出一页可编辑 PPT。\n\n"
                    "2. 多页多 Agent 高质量模式\n"
                    "适合：你想转换整份 PPTX、质量优先，并且你能接受非常高的 token 消耗。\n"
                    "建议：仅限 Codex Pro 订阅用户使用。\n"
                    "消耗：非常高；每页会交给独立转换任务处理。\n"
                    "结果：输出完整可编辑 PPT；失败页会保留原图并加提示。\n\n"
                    "如果你更想省 token，建议先选单页省 token 模式。"
                    "如果要转换整份 PPTX，则使用多页多 Agent 高质量模式；它非常消耗 token，"
                    "建议仅限 Codex Pro 订阅用户使用。\n\n"
                    "请告诉我你的选择。"
                ),
                "options": [
                    "single-page-token-saving",
                    "multi-agent-high-quality",
                ],
                "token_economy_note": (
                    "For the lowest token cost, recommend single-page-token-saving first. "
                    "Full-deck multi-agent conversion costs more because each selected page "
                    "is handled as an independent conversion task."
                ),
                "internal_recording_note": (
                    "Record the user's answer as conversion_mode_gate.user_choice. "
                    "Use single-page-token-saving for one PNG or one sample slide. "
                    "Use multi-agent-high-quality when the user chooses full-deck PPTX conversion "
                    "with independent per-slide workers."
                ),
            },
            "ocr_runtime_gate": {
                "requires_user_confirmation": not ocr_ready,
                "needs_setup": not ocr_ready,
                "auto_pass_existing_runtime": ocr_ready,
                "explain_to_user": (
                    "OCR reads each source slide image and provides source text, text "
                    "positions, line breaks, and Chinese text usability checks. It is "
                    "required for reliable text removal, editable text placement, and QA."
                ),
                "why_setup_takes_time": (
                    None
                    if ocr_ready
                    else (
                        "OCR runtime is not ready yet. First-time setup may take several "
                        "minutes because Python dependencies and OCR model files may need "
                        "to be installed or downloaded."
                    )
                ),
                "setup_command": (
                    None
                    if ocr_ready
                    else f"python \"{setup_script}\" --gates-file GATES_FILE --yes"
                ),
            },
            "scope_and_worker_gate": {
                "question": (
                    f"这份 PPT 一共有 {slide_count} 页。\n\n"
                    "你想怎么转换？\n"
                    f"1. 转换全部 {slide_count} 页\n"
                    "2. 只转换指定页\n"
                    "3. 先转换 1 页看看效果\n\n"
                    "说明：转换页数越多，消耗的 token 越多；多页多 Agent 高质量模式非常消耗 token，"
                    "建议仅限 Codex Pro 订阅用户使用。"
                    "如果你选择多页多 Agent 高质量模式，我会把每一页交给独立转换任务处理；"
                    "这些任务会读取对应页面的图片和文字内容，token 消耗更高。\n\n"
                    "请告诉我你的选择。"
                ),
                "options": list(SCOPE_GATE_CHOICES),
                "worker_modes": list(WORKER_MODE_CHOICES),
                "pro_subscription_note": (
                    "Multi-agent high-quality mode is better suited to Pro subscribers "
                    "because it may start one independent worker per selected page and "
                    "therefore consumes more tokens."
                ),
                "recommendation": (
                    "Ask this plain-language question verbatim or translate it naturally. "
                    "Do not ask the user to approve app-native/external-codex/worker_mode "
                    "fields directly. If the user is unsure, recommend one sample page first."
                ),
                "internal_recording_note": (
                    "After the user chooses, map the answer into scope_and_worker_gate. "
                    "For normal in-app multi-agent execution, record worker_mode=app-native, "
                    "worker_execution_approved=true, token_cost_acknowledged=true, and "
                    "content_sharing_acknowledged=true."
                ),
            },
        },
        "next_steps": [
            "If OCR is not passed-text-usable, ask the OCR Runtime Gate before setup or limited fallback. If OCR is already passed-text-usable, record auto_passed_existing_runtime=true.",
            "Ask the plain-language page-scope question from user_confirmation_required before creating slide jobs; do not expose internal worker_mode field names to the user.",
            "Write the user's answers into a gates file using gates_file_template.",
            "If the user confirmed OCR setup, run setup_ocr_runtime.py --gates-file GATES_FILE --yes, then check_ocr_runtime.py --json.",
            "Initialize: deck_controller.py SOURCE --run-dir RUN_DIR --gates-file GATES_FILE --all-slides|--slides \"...\" --ocr-python OCR_PYTHON",
        ],
        "gates_file_template": gates_file_template(),
    }


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def write_json(path: Path, data: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def read_attempt_json(attempt_dir: Path, name: str) -> dict | None:
    path = attempt_dir / name
    if not path.exists():
        return None
    return read_json(path)


def report_summary(report: dict | None) -> dict:
    if not isinstance(report, dict):
        return {}
    summary = report.get("summary")
    return summary if isinstance(summary, dict) else report


def validate_region_crop_plan_contract(summary: dict) -> tuple[bool, str | None]:
    issues = region_crop_plan_issues(summary.get("region_crop_plan"))
    if not issues:
        return True, None
    issue_names = []
    for issue in issues[:5]:
        name = str(issue.get("issue", "region_crop_plan_issue"))
        region_id = issue.get("region_id")
        issue_names.append(f"{name}:{region_id}" if region_id else name)
    return False, "region_crop_plan_contract_failed: " + "; ".join(issue_names)


def validate_visual_strategy_artifacts(attempt_dir: Path, summary: dict) -> tuple[bool, str | None]:
    issues: list[str] = []

    try:
        plan = read_attempt_json(attempt_dir, "source_reconstruction_plan.json")
    except json.JSONDecodeError as exc:
        return False, f"source_reconstruction_plan_invalid_json: {exc}"
    if isinstance(plan, dict):
        native_issues = high_risk_native_strategy_issues(plan)
        for issue in native_issues[:5]:
            region_id = issue.get("region_id") or "unknown-region"
            visual_type = issue.get("visual_type") or "unknown-visual-type"
            actual = issue.get("actual_strategy") or issue.get("expected_strategy")
            issues.append(f"high_risk_native_strategy:{region_id}:{visual_type}:{actual}")

        try:
            packaging_report = read_attempt_json(attempt_dir, "packaging_report.json")
        except json.JSONDecodeError as exc:
            return False, f"packaging_report_invalid_json: {exc}"
        zero_crop_issue = zero_crop_high_native_issue(plan, report_summary(packaging_report))
        if zero_crop_issue:
            issues.append(
                "high_risk_zero_source_crops_all_native:"
                f"pictures={zero_crop_issue.get('pictures')};"
                f"croppedPictures={zero_crop_issue.get('croppedPictures')};"
                f"nativeShapes={zero_crop_issue.get('nativeShapes')};"
                f"nativeLines={zero_crop_issue.get('nativeLines')}"
            )
        full_slide_issue = full_slide_textless_background_issue(
            plan,
            report_summary(packaging_report),
            summary,
        )
        if full_slide_issue:
            issues.append(
                "full_slide_textless_background_disguised_as_region_crops:"
                f"pictures={full_slide_issue.get('pictures')};"
                f"croppedPictures={full_slide_issue.get('croppedPictures')};"
                f"regions={full_slide_issue.get('disguised_region_count')}/"
                f"{full_slide_issue.get('region_count')};"
                f"fullSlidePictures={','.join(full_slide_issue.get('full_slide_picture_ids') or [])}"
            )

    try:
        visual_report = read_attempt_json(attempt_dir, "visual_strategy_report.json")
    except json.JSONDecodeError as exc:
        return False, f"visual_strategy_report_invalid_json: {exc}"
    if isinstance(visual_report, dict):
        warnings = visual_report.get("warnings", [])
        if isinstance(warnings, list):
            for warning in warnings:
                if not isinstance(warning, dict):
                    continue
                if warning.get("warning") == "native_redraw_used_for_high_risk_region":
                    issues.append(
                        "high_risk_native_strategy_warning:"
                        f"{warning.get('id') or warning.get('region_id') or 'unknown-region'}"
                    )
                if warning.get("warning") == "full_slide_textless_background_disguised_as_region_crops":
                    issues.append(
                        "full_slide_textless_strategy_warning:"
                        f"{warning.get('id') or warning.get('region_id') or 'strategy_review'}"
                    )

    if issues:
        return False, "; ".join(issues[:8])
    return True, None


def utf8_subprocess_env() -> dict[str, str]:
    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"
    return env


def resolve_rel_target(part_name: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(part_name), target))


def read_slide_size(zf: zipfile.ZipFile) -> tuple[int, int]:
    if "ppt/presentation.xml" not in zf.namelist():
        raise ValueError("ppt/presentation.xml missing")
    root = ET.fromstring(zf.read("ppt/presentation.xml"))
    size = root.find("p:sldSz", NS)
    if size is None:
        raise ValueError("presentation slide size missing")
    return int(size.attrib["cx"]), int(size.attrib["cy"])


def _ocr_import_check_script() -> str:
    return r"""
import json
import os
import sys

deps = sys.argv[1] if len(sys.argv) > 1 else ""
added_paths = []
if deps:
    for rel in ("shapely.libs", "numpy.libs", os.path.join("onnxruntime", "capi")):
        path = os.path.join(deps, rel)
        if os.path.isdir(path):
            added_paths.append(path)
            try:
                os.add_dll_directory(path)
            except Exception:
                pass
    if os.path.isdir(deps):
        added_paths.append(deps)
        sys.path.insert(0, deps)

checks = {}
for module in ("shapely.geometry", "rapidocr_onnxruntime"):
    try:
        imported = __import__(module)
    except Exception as exc:
        checks[module] = {"status": "unavailable", "error": repr(exc)}
    else:
        checks[module] = {
            "status": "available",
            "error": None,
            "file": getattr(imported, "__file__", None),
        }

status = "passed" if all(v["status"] == "available" for v in checks.values()) else "warning"
print(json.dumps({"status": status, "checks": checks, "added_paths": added_paths}, ensure_ascii=False))
"""


def normalize_ocr_runtime_status(value: str | None, dependency_status: str | None = None) -> str:
    if value in {OCR_TEXT_USABLE_STATUS, OCR_GEOMETRY_ONLY_STATUS, OCR_FAILED_STATUS}:
        return value
    if dependency_status == "passed":
        return OCR_GEOMETRY_ONLY_STATUS
    return OCR_FAILED_STATUS


def ocr_text_is_usable(value: str | None) -> bool:
    return value == OCR_TEXT_USABLE_STATUS


def get_ocr_dependency_status(
    ocr_python: Path | None = None,
    ocr_deps: Path | None = None,
    check_dir: Path | None = None,
) -> dict:
    runtime = {
        "mode": "external" if ocr_python else "current-python",
        "python": str(ocr_python) if ocr_python else sys.executable,
        "deps": str(ocr_deps) if ocr_deps else None,
    }
    checker = Path(__file__).with_name("check_ocr_runtime.py")
    if checker.exists():
        python_exe = str(ocr_python) if ocr_python else sys.executable
        command = [python_exe, str(checker), "--json"]
        if ocr_deps:
            command.extend(["--deps", str(ocr_deps)])
        if check_dir:
            command.extend(["--output-dir", str(check_dir)])
        try:
            proc = subprocess.run(
                command,
                check=False,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                env=utf8_subprocess_env(),
            )
        except Exception as exc:  # noqa: BLE001 - preflight should record runtime launch failures
            return {
                "status": "warning",
                "ocr_runtime_status": OCR_FAILED_STATUS,
                "runtime": runtime,
                "checks": {},
                "error": repr(exc),
                "checker": str(checker),
            }
        if proc.returncode == 0:
            try:
                result = json.loads(proc.stdout)
            except json.JSONDecodeError:
                result = None
            if isinstance(result, dict):
                checker_status = result.get("ocr_runtime_status") or result.get("status")
                dependency_status = result.get("dependency_status") or (
                    "passed" if checker_status in {OCR_TEXT_USABLE_STATUS, OCR_GEOMETRY_ONLY_STATUS} else "warning"
                )
                result["status"] = dependency_status
                result["ocr_runtime_status"] = normalize_ocr_runtime_status(
                    checker_status,
                    dependency_status,
                )
                result["runtime"] = runtime | result.get("runtime", {})
                result["checker"] = str(checker)
                return result
        return {
            "status": "warning",
            "ocr_runtime_status": OCR_FAILED_STATUS,
            "runtime": runtime,
            "checks": {},
            "error": proc.stderr.strip() or proc.stdout.strip() or f"checker_returncode={proc.returncode}",
            "checker": str(checker),
        }

    if ocr_python:
        if not ocr_python.exists():
            return {
                "status": "warning",
                "ocr_runtime_status": OCR_FAILED_STATUS,
                "runtime": runtime,
                "checks": {},
                "error": f"ocr_python_not_found: {ocr_python}",
            }
        try:
            proc = subprocess.run(
                [str(ocr_python), "-c", _ocr_import_check_script(), str(ocr_deps or "")],
                check=False,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                env=utf8_subprocess_env(),
            )
        except Exception as exc:  # noqa: BLE001 - preflight should capture runtime launch failures
            return {
                "status": "warning",
                "ocr_runtime_status": OCR_FAILED_STATUS,
                "runtime": runtime,
                "checks": {},
                "error": repr(exc),
            }
        if proc.returncode != 0:
            return {
                "status": "warning",
                "ocr_runtime_status": OCR_FAILED_STATUS,
                "runtime": runtime,
                "checks": {},
                "error": proc.stderr.strip() or proc.stdout.strip(),
            }
        try:
            result = json.loads(proc.stdout)
        except json.JSONDecodeError as exc:
            return {
                "status": "warning",
                "ocr_runtime_status": OCR_FAILED_STATUS,
                "runtime": runtime,
                "checks": {},
                "error": f"invalid_ocr_preflight_json: {exc}: {proc.stdout}",
            }
        result["runtime"] = runtime
        result["ocr_runtime_status"] = normalize_ocr_runtime_status(None, result.get("status"))
        return result

    checks = {}
    for module in ("shapely.geometry", "rapidocr_onnxruntime"):
        try:
            __import__(module)
        except Exception as exc:  # noqa: BLE001 - preflight should record exact import problem
            checks[module] = {"status": "unavailable", "error": repr(exc)}
        else:
            checks[module] = {"status": "available", "error": None}

    status = "passed" if all(v["status"] == "available" for v in checks.values()) else "warning"
    return {
        "status": status,
        "ocr_runtime_status": normalize_ocr_runtime_status(None, status),
        "runtime": runtime,
        "checks": checks,
    }


def extract_image_only_slides(source_pptx: Path, run_dir: Path) -> tuple[dict, list[dict]]:
    source_dir = run_dir / "source_slides"
    source_dir.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(source_pptx) as zf:
        slide_width_emu, slide_height_emu = read_slide_size(zf)
        slide_names = sorted(
            [name for name in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)],
            key=slide_sort_key,
        )
        if not slide_names:
            raise ValueError("No slides found in PPTX")

        slide_infos: list[dict] = []
        first_image_size: tuple[int | None, int | None] = (None, None)

        for idx, slide_name in enumerate(slide_names, start=1):
            slide_xml = ET.fromstring(zf.read(slide_name))
            rels_name = f"ppt/slides/_rels/slide{idx}.xml.rels"
            rel_map = {}
            if rels_name in zf.namelist():
                rels_xml = ET.fromstring(zf.read(rels_name))
                for rel in rels_xml.findall("rel:Relationship", NS):
                    rid = rel.attrib.get("Id")
                    target = rel.attrib.get("Target")
                    rel_type = rel.attrib.get("Type")
                    if rid and target:
                        rel_map[rid] = {
                            "target": resolve_rel_target(slide_name, target),
                            "type": rel_type,
                        }

            pics = []
            for pic_i, pic in enumerate(slide_xml.findall(".//p:pic", NS), start=1):
                blip = pic.find(".//a:blip", NS)
                rid = blip.attrib.get(f"{{{NS['r']}}}embed") if blip is not None else None
                xfrm = pic.find(".//a:xfrm", NS)
                off = xfrm.find("a:off", NS) if xfrm is not None else None
                ext = xfrm.find("a:ext", NS) if xfrm is not None else None
                target = rel_map.get(rid, {}).get("target") if rid else None
                extracted = None
                image_size = None
                image_format = None

                if target and target in zf.namelist():
                    suffix = Path(target).suffix or ".bin"
                    extracted = source_dir / f"slide_{idx:02d}{suffix.lower()}"
                    with zf.open(target) as src, extracted.open("wb") as dst:
                        shutil.copyfileobj(src, dst)
                    with Image.open(extracted) as img:
                        image_size = {"width_px": img.width, "height_px": img.height}
                        image_format = img.format
                        if first_image_size == (None, None):
                            first_image_size = (img.width, img.height)

                x = int(off.attrib.get("x", 0)) if off is not None else None
                y = int(off.attrib.get("y", 0)) if off is not None else None
                cx = int(ext.attrib.get("cx", 0)) if ext is not None else None
                cy = int(ext.attrib.get("cy", 0)) if ext is not None else None
                is_full_slide = (
                    x == 0
                    and y == 0
                    and cx == slide_width_emu
                    and cy == slide_height_emu
                )

                pics.append(
                    {
                        "picture_index": pic_i,
                        "relationship_id": rid,
                        "media_target": target,
                        "extracted_path": str(extracted) if extracted else None,
                        "image_format": image_format,
                        "image_size": image_size,
                        "x": x,
                        "y": y,
                        "cx": cx,
                        "cy": cy,
                        "is_full_slide": is_full_slide,
                    }
                )

            slide_infos.append(
                {
                    "slide_index": idx,
                    "slide_xml": slide_name,
                    "picture_count": len(pics),
                    "pictures": pics,
                    "is_single_full_slide_picture": (
                        len(pics) == 1 and pics[0]["is_full_slide"] and pics[0]["extracted_path"]
                    ),
                }
            )

    width_px, height_px = first_image_size
    slide_size = {
        "width_emu": slide_width_emu,
        "height_emu": slide_height_emu,
        "width_px": width_px,
        "height_px": height_px,
        "aspect_ratio": "16:9" if slide_width_emu * 9 == slide_height_emu * 16 else None,
    }
    return slide_size, slide_infos


def make_slide_job_manifest(
    run_id: str,
    slide: dict,
    run_dir: Path,
    ocr_python: Path | None = None,
    ocr_deps: Path | None = None,
    ocr_runtime_status: str = OCR_FAILED_STATUS,
    ocr_runtime_report: dict | None = None,
    worker_mode: str = "app-native",
    conversion_mode: str = "multi-agent-high-quality",
) -> dict:
    if conversion_mode in RETIRED_CONVERSION_MODE_CHOICES or worker_mode in RETIRED_WORKER_MODE_CHOICES:
        raise ValueError(
            "Sequential token-saving execution is retired and not supported in this release. "
            "Use multi-agent-high-quality with a real per-slide worker, or single-page-token-saving "
            "for one PNG/sample slide."
        )
    picture = slide["pictures"][0]
    slide_dir = run_dir / "slide_jobs" / f"slide_{slide['slide_index']:02d}"
    attempt_dir = slide_dir / "attempt_01"
    attempt_dir.mkdir(parents=True, exist_ok=True)

    source_image = Path(picture["extracted_path"])
    source_size = picture["image_size"] or {"width_px": None, "height_px": None}
    page_contract = make_page_conversion_contract(slide["slide_index"], source_size)
    crop_contract = make_crop_manifest_contract(slide["slide_index"], source_size)
    page_contract_path = slide_dir / "page_conversion_contract.json"
    crop_contract_path = slide_dir / "crop_manifest_contract.json"
    reference_paths = context_reference_paths()
    compact_protocol_path = ensure_worker_compact_protocol(run_dir)
    required_reading = [
        {
            "id": "worker_compact_protocol",
            "path": str(compact_protocol_path),
            "must_read": True,
            "reason": "Short v3 preview run-specific worker protocol. Read this instead of the full long reference set by default.",
        },
        {
            "id": "slide_job_manifest",
            "path": str(slide_dir / "slide_job_manifest.json"),
            "must_read": True,
            "reason": "Per-slide input, output, OCR runtime, and status contract.",
        },
        {
            "id": "page_conversion_contract",
            "path": str(page_contract_path),
            "must_read": True,
            "reason": "Short per-slide execution lock for fonts, effects, text backgrounds, crop-first judgment, and fidelity tradeoffs.",
        },
        {
            "id": "crop_manifest_contract",
            "path": str(crop_contract_path),
            "must_read": True,
            "reason": "Short per-slide crop planning contract and required crop_manifest.json review gates.",
        },
    ]
    fallback_reading = fallback_reading_entries(reference_paths)
    sequential_mode = False
    execution_model = "delegated-real-per-slide-worker"
    allowed_execution_models = [execution_model]
    manifest = {
        "schema_version": "ppt-to-editable-slide-job-v1",
        "deck_run_id": run_id,
        "slide_index": slide["slide_index"],
        "source_image": str(source_image),
        "source_slide_size": source_size,
        "output_dir": str(attempt_dir),
        "attempt": 1,
        "max_attempts": 2,
        "worker_contract": {
            "execution_model": execution_model,
            "reuse_skill": PRIMARY_SINGLE_PAGE_ENGINE_ID,
            "primary_single_page_engine": {
                "id": PRIMARY_SINGLE_PAGE_ENGINE_ID,
                "root": str(single_page_engine_root()),
                "skill": reference_paths["single_page_engine_skill_3_1"],
                "script_paths": single_page_engine_script_paths(),
                "active_skill_policy": (
                    "Do not use an external local development folder as the primary engine. "
                    "Use this v3.1 preview package and mark the slide failed-retryable if its "
                    "required scripts are unreadable."
                ),
            },
            "task_scope": "single-slide-png-to-editable-pptx",
            "must_not_process_other_slides": True,
            "must_run_in_separate_worker_context": not sequential_mode,
            "controller_agent_may_process_slides_sequentially": sequential_mode,
            "must_report_worker_identity": True,
            "allowed_execution_models": allowed_execution_models,
            "worker_identity_required": [
                "worker_thread_id or worker_agent_id",
                "execution_model",
                "slide_scope",
            ],
            "controller_agent_must_not_generate_slide_pptx": not sequential_mode,
            "forbidden_batch_pattern": (
                "A controller or test thread must not create one script that builds multiple slide.pptx outputs."
                if not sequential_mode
                else "Sequential mode may process listed slide jobs one at a time, but must not create one unreviewed batch builder that bypasses per-slide QA."
            ),
        },
        "required_reading": required_reading,
        "fallback_reading": fallback_reading,
        "ocr_runtime": {
            "mode": "external" if ocr_python else "current-python",
            "python": str(ocr_python) if ocr_python else None,
            "deps": str(ocr_deps) if ocr_deps else None,
            "status": normalize_ocr_runtime_status(ocr_runtime_status),
            "report": ocr_runtime_report or {},
            "development_note": (
                "Use this package's setup_ocr_runtime and check_ocr_runtime flow. "
                "Do not rely on a local development OCR virtual environment."
            )
            if ocr_python
            else None,
        },
        "required_outputs": [
            "progress.json",
            "crop_manifest.json",
            "slide.pptx",
            "render_report.json",
            "visual-qa-comparison.html",
            "editability_report.json",
            "qa_summary.json",
        ],
        "page_conversion_contract": page_contract,
        "crop_manifest_contract": crop_contract,
        "routing_decision_table": routing_decision_table(),
        "reference_loading_policy": reference_loading_policy(required_reading, fallback_reading),
        "quality_rules": {
            "meaningful_text_should_be_editable": True,
            "multi_line_text_should_remain_grouped": True,
            "multi_line_grouping_is_conservative_same_style_only": True,
            "complex_visuals_can_use_tight_crop": True,
            "complex_shapes_should_use_textless_crop_when_native_redraw_would_drift": True,
            "must_follow_page_conversion_contract": True,
            "must_write_crop_manifest_before_packaging_when_any_crop_or_complex_region_is_used": True,
            "must_follow_routing_decision_table_before_reconstruction_plan": True,
            "must_minimize_reference_loading": True,
            "must_not_read_historical_round_outputs_without_user_request": True,
            "do_not_fake_reconstruction_with_full_slide_image": True,
            "must_read_required_references": True,
            "must_report_references_read": True,
            "per_slide_output_is_draft_until_deck_level_qa": True,
            "must_apply_microsoft_yahei_font_policy": True,
            "must_preserve_relative_font_size_consistency": True,
            "must_not_shrink_text_to_fit_when_source_scale_would_be_lost": True,
            "must_clear_unwanted_shape_effects": True,
            "must_review_powerpoint_rendered_png_when_available": True,
            "must_plan_region_level_textless_crops_for_complex_pages": True,
            "must_not_use_whole_slide_textless_crop_without_explicit_acceptance": True,
            "must_use_crop_first_for_complex_artwork_regions": True,
            "must_not_add_invented_text_background_fills": True,
            "must_preserve_non_text_visual_anchors": True,
            "must_not_erase_arrows_connectors_step_markers_or_card_edges": True,
            "must_write_full_progress_heartbeat": True,
            "must_report_real_worker_identity": True,
            "must_use_single_page_engine_3_2": True,
            "must_report_single_page_engine_used": True,
            "must_report_single_page_engine_scripts_used": True,
            "must_report_fallback_references_read": True,
            "must_prefer_integrated_deterministic_packager": True,
            "must_not_write_full_ad_hoc_packager_when_integrated_packager_applies": True,
            "ocr_unavailable_blocks_success_status": True,
            "ocr_text_unusable_blocks_success_status": True,
            "accepted_with_limitations_is_minor_only": True,
            "blocking_limitations_trigger_retry_or_fallback": BLOCKING_SUCCESS_LIMITATION_PATTERNS,
        },
        "failure_writeback": {
            "status_values": [
                "passed",
                "needs-review",
                "failed-retryable",
                "image-fallback",
                "accepted-with-limitations",
            ],
            "required_on_failure": [
                "failure_reason",
                "failed_step",
                "recommended_retry_change",
            ],
        },
        "qa_summary_required_fields": [
            "status",
            "mode",
            "references_read",
            "fallback_references_read",
            "font_policy_applied",
            "font_size_consistency_review",
            "effects_cleared",
            "rendered_png_reviewed",
            "render_report",
            "region_crop_plan",
            "crop_manifest",
            "complex_visual_strategy_review",
            "text_background_policy",
            "non_text_anchor_preservation_review",
            "editable_text_count",
            "native_shapes_count",
            "source_crop_count",
            "visual_qa",
            "failure_reason",
            "limitations",
            "recommended_next_round_change",
            "worker_identity",
            "ocr_runtime_status",
            "single_page_engine_used",
            "single_page_engine_scripts_used",
        ],
    }
    write_json(page_contract_path, page_contract)
    write_json(crop_contract_path, crop_contract)
    write_json(slide_dir / "slide_job_manifest.json", manifest)
    write_agent_task(slide_dir / "AGENT_TASK.md", manifest)
    return manifest


def write_agent_task(path: Path, manifest: dict) -> None:
    required_reading = "\n".join(
        f"- `{item['path']}`  # {item['reason']}" for item in manifest["required_reading"]
    )
    fallback_reading = "\n".join(
        f"- `{item['path']}`  # {item['reason']}" for item in manifest.get("fallback_reading", [])
    )
    required_reading_ids = ", ".join(f"`{item['id']}`" for item in manifest["required_reading"])
    fallback_reading_ids = ", ".join(
        f"`{item['id']}`" for item in manifest.get("fallback_reading", [])
    )
    required_reading_paths_by_id = {
        item["id"]: item["path"] for item in manifest["required_reading"]
    }
    handoff_script = str(Path(__file__).with_name("write_slide_qa_summary_from_single_page.py"))
    engine = manifest["worker_contract"]["primary_single_page_engine"]
    engine_scripts = "\n".join(
        f"- `{name}`: `{script_path}`"
        for name, script_path in engine["script_paths"].items()
    )
    routing_table = json.dumps(manifest["routing_decision_table"], ensure_ascii=False, indent=2)
    reference_policy = json.dumps(manifest["reference_loading_policy"], ensure_ascii=False, indent=2)
    worker_contract = manifest["worker_contract"]
    execution_model = worker_contract.get("execution_model", "delegated-real-per-slide-worker")
    execution_contract_line = (
        "Run this task in a real separate per-slide worker context. "
        "The controller/test thread must not build this slide itself."
    )
    worker_id_placeholder = "<your real worker_agent_id>"
    text = f"""# Slide {manifest['slide_index']:02d} Agent Task

Use the bundled `ppt-to-editable` v3.1 preview single-page engine on exactly one source image.

## Input

```text
{manifest['source_image']}
```

## Required Contract

- Read the compact **Required Reading** set before conversion work.
- Do not read **Fallback Reading** files unless blocked, a QA gate fails, or the compact protocol is ambiguous.
- {execution_contract_line}
- Report your real worker identity in `qa_summary.json` as `worker_identity`.
- Write `progress.json` at every required milestone: `started`, `ocr_done`, `plan_done`, `packaging_done`, `qa_done`.
- Do not process other slides.
- Use this `ppt-to-editable` v3.1 preview package as the primary single-page conversion engine.
- Do not use an external local development folder as the primary engine for this run.
- Read `page_conversion_contract.json` and `crop_manifest_contract.json` before writing a reconstruction plan.
- Classify regions with the Routing Decision Table before deciding what to rebuild or crop.
- Do not read historical test reports, old round outputs, prior visual QA HTML, or baseline artifacts unless the user explicitly asks to reuse them for diagnosis.
- Prefer the integrated deterministic reconstruction packager and QA scripts over writing a full one-off PPTX builder.
- Do not force native reconstruction for complex shapes, bespoke arrows, icon groups, shaded cards, or source-specific modules. If native redraw would drift or look unnatural, use a tight local crop or textless crop plus editable text for that region.
- Do not shrink body text into tiny outliers just to make a manually guessed layout fit. Preserve the source slide's relative font hierarchy; widen or split regions before reducing readable body text below the source scale.
- Use crop-first judgment for complex artwork regions: product photos, detailed icons, curved paths, gradient arrows, matrix backgrounds, textured modules, and source-specific card stacks should be source-derived crops or textless crops.
- In v3.1, high-risk regions must not use `native_redraw`, `native_lines`, `native_shapes`, or `native_reconstruction`; `override_reason` does not bypass this gate.
- Every `region_crop_plan` entry must be an object with `id`, `strategy`, `source_bbox_px`, and `reason`; do not use strings or `not-reported` placeholders.
- Do not add pale rectangles or invented backing fills behind editable text. Text boxes should be transparent unless the source itself has a visible card/bar/background there; repair old source text with a textless source crop instead of covering it with a generic fill.
- Do not erase non-text visual anchors. Preserve arrows, arrowheads, route lines, connectors, step markers, card edges, cut corners, dividers, and icons even when removing source text from a crop.
- If text removal damages a required anchor, enlarge/split the crop, keep the original anchor as a source crop, or redo the local textless crop before returning success.
- Route this slide independently.
- Treat your output as a draft until deck-level QA accepts it.
- Write outputs under:

```text
{manifest['output_dir']}
```

## Required Reading

{required_reading}

If any required compact file is missing or unreadable, stop and write `qa_summary.json` with `status=failed-retryable`.

## Fallback Reading

Read these longer references only when blocked, when a QA gate fails, or when the compact protocol is ambiguous. If you read any of them, report the ids and short reasons in `fallback_references_read`.

{fallback_reading}

## Page Execution Contracts

Read these two small JSON contracts first. They are the local execution lock for this slide:

```text
{required_reading_paths_by_id['page_conversion_contract']}
{required_reading_paths_by_id['crop_manifest_contract']}
```

`page_conversion_contract.json` owns default font, effect, text-background, and fidelity tradeoff rules. `crop_manifest_contract.json` owns crop planning and textless-crop review gates.

## Routing Decision Table

Use this table before writing `source_reconstruction_plan.json`:

```json
{routing_table}
```

## Reference Loading Policy

Keep context small:

```json
{reference_policy}
```

## Primary Single-Page Engine

```text
{engine['root']}
```

Required script paths:

{engine_scripts}

Execution rule:

- If the slide is reconstruction or mixed reconstruction, write/review `source_reconstruction_plan.json`, then run `run_reconstruction_qa.py` or `package_reconstruction_deck.py` from this v3.1 preview package.
- If the slide is clean-background, use the integrated clean-background packager path and report it.
- If you write small helper scripts, keep them limited to plan/crop preparation. Do not replace the integrated packager with a full custom builder unless the integrated packager is proven inapplicable; in that case mark the slide `failed-retryable` and explain why.
- In `qa_summary.json`, set `single_page_engine_used` to `{PRIMARY_SINGLE_PAGE_ENGINE_ID}` and list every integrated v3.1 script actually called in `single_page_engine_scripts_used`.
- For complex shapes, process arrows, icon clusters, shaded cards, or detailed visual modules, prefer textless crop plus editable text when native shape reconstruction would reduce fidelity.

## Required Outputs

- `progress.json`
- `slide.pptx`
- `render_report.json`
- `visual-qa-comparison.html`
- `editability_report.json`
- `qa_summary.json`

## Required QA Summary Fields

`qa_summary.json` must include:

- `status`
- `mode`
- `references_read`
- `fallback_references_read`
- `font_policy_applied`
- `font_size_consistency_review`
- `effects_cleared`
- `rendered_png_reviewed`
- `render_report`
- `region_crop_plan`
- `crop_manifest`
- `complex_visual_strategy_review`
- `text_background_policy`
- `non_text_anchor_preservation_review`
- `editable_text_count`
- `native_shapes_count`
- `source_crop_count`
- `visual_qa`
- `failure_reason`
- `limitations`
- `recommended_next_round_change`
- `worker_identity`
- `ocr_runtime_status`
- `single_page_engine_used`
- `single_page_engine_scripts_used`

`references_read` must include these ids: {required_reading_ids}.
`fallback_references_read` should be an empty list when no fallback reference was needed. If any fallback reference was read, use these ids only: {fallback_reading_ids}.

`worker_identity` must include at least one real worker id:

```json
{{
  "worker_thread_id": "Codex thread id when this slide was handled in a separate thread",
  "worker_agent_id": "agent id if the runtime exposes one",
  "execution_model": "{execution_model}",
  "slide_scope": [{manifest['slide_index']}]
}}
```

Do not invent different ids inside one batch script. If the controller/test thread generated multiple slide PPTX files itself, this slide must be marked `failed-retryable`.

The controller dispatched this slide with OCR runtime status:

```text
{manifest['ocr_runtime']['status']}
```

Copy this value into `qa_summary.json` as `ocr_runtime_status` unless your slide-local OCR check proves a stricter failure. The only allowed values are `passed-text-usable`, `passed-geometry-only`, and `failed`.

If OCR is not `passed-text-usable`, do not mark the slide `passed` or `accepted-with-limitations`. Use `failed-retryable` for a repairable conversion or `image-fallback` for a deliberate fallback.

If PowerPoint or LibreOffice is available, export the generated `slide.pptx` to a rendered PNG before returning the slide. Use that rendered PNG to fix text overlap, text overflow, font fallback, unwanted shadows/effects, or region-crop drift. Report the result in `rendered_png_reviewed`.

Use the bundled render helper when possible:

```text
python "{Path(__file__).with_name('render_pptx_qa.py')}" "{manifest['output_dir']}\\slide.pptx" --out-dir "{manifest['output_dir']}\\render" --source "{manifest['source_image']}" --report "{manifest['output_dir']}\\render_report.json"
```

Report the render report path in `qa_summary.json` as `render_report`. A success status requires `render_report.json` with `status=rendered`.

Use Microsoft YaHei for all generated text runs, clear unwanted native shape effects, and record these in `font_policy_applied` and `effects_cleared`.

Check font size consistency against the source and record it in `font_size_consistency_review`. A success status is not allowed when visible body/callout text becomes a tiny outlier relative to the source.

Group same-style multiline body/callout text into one editable text box when font face, size, weight, color, alignment, column/region, and spacing match. Do not merge titles with body text, labels, table cells, legends, axis labels, badge numbers, step labels, mixed-style text, or cross-column text. If a same-style split must remain line-level, mark it with `line_level_trace_required` or `do_not_merge`.

For complex structured pages, record a semantic 3-5 region crop/native/text plan in `region_crop_plan`; do not use one whole-slide textless crop unless explicitly accepted as fallback. If multiple regions all point to the same whole-slide/background picture, mark the slide `failed-retryable` or `image-fallback`; do not mark it `passed` or `accepted-with-limitations`.

Write `crop_manifest.json` before packaging whenever the page has source crops, textless crops, clean badge composites, or complex visual regions. If the page truly needs no crops, still write an empty manifest with a reason such as `no_source_crops_required`.

Record `complex_visual_strategy_review` and explain which complex regions used crop-first handling. A success status is not allowed when complex artwork was redrawn natively and became visibly worse, or when the claimed region crops are actually one full-slide visual layer.

Record `text_background_policy`. A success status is not allowed when editable text was made readable by adding invented pale rectangles, generic backing fills, or non-source-matched text backgrounds.

Record `non_text_anchor_preservation_review`. A success status is not allowed when arrows, arrowheads, route lines, connectors, step markers, card edges, cut corners, dividers, or icons are missing, broken, cropped, or erased.

`accepted-with-limitations` is only for minor visual tradeoffs such as non-editable icons or slight texture mismatch. Do not use it for OCR unavailable, manual layout guessing, major layout drift, rendered text overflow, font fallback, or obvious text-size drift. Those are `failed-retryable` unless this attempt is explicitly being converted to `image-fallback`.

After the single-page QA chain writes its own `qa_summary.json`, run this v3.1 handoff command to convert it into the controller-required summary. Replace the worker id value with your real agent id or thread id; do not invent one:

```text
python "{handoff_script}" --slide-job-manifest "{Path(manifest['output_dir']).parent / 'slide_job_manifest.json'}" --attempt-dir "{manifest['output_dir']}" --worker-agent-id "{worker_id_placeholder}" --output "{manifest['output_dir']}\\qa_summary.json"
```

This command preserves the original single-page summary as `qa_summary.single-page.json`, writes the full v3.1 `qa_summary.json`, normalizes render reports for the v3.1 controller, and appends `qa_done` to `progress.json`.

If the slide fails, write `qa_summary.json` with `failure_reason`, `failed_step`, and `recommended_retry_change`.
"""
    path.write_text(text, encoding="utf-8")


def build_deck_manifest(
    source_pptx: Path,
    run_dir: Path,
    run_id: str,
    slide_size: dict,
    slides: list[dict],
    source_slide_count: int | None = None,
    ocr_python: Path | None = None,
    ocr_deps: Path | None = None,
    allow_ocr_not_ready: bool = False,
    gates: dict | None = None,
    gates_file: Path | None = None,
) -> dict:
    ocr_status = get_ocr_dependency_status(ocr_python, ocr_deps, run_dir / "ocr_runtime_check")
    ocr_runtime_status = normalize_ocr_runtime_status(
        ocr_status.get("ocr_runtime_status"),
        ocr_status.get("status"),
    )
    require_ocr_text_usable(ocr_runtime_status, allow_ocr_not_ready)
    all_image_only = all(s["is_single_full_slide_picture"] for s in slides)
    all_png = all(
        (s["pictures"][0].get("image_format") == "PNG") for s in slides if s["pictures"]
    )
    notes = []
    if ocr_status["status"] != "passed":
        notes.append(
            "OCR dependency preflight is warning; deck controller can still create manifests and fallback deck, "
            "but quality editable outputs must not be accepted as passed or accepted-with-limitations."
        )
    if ocr_runtime_status != OCR_TEXT_USABLE_STATUS:
        notes.append(
            "OCR runtime is not passed-text-usable; single-slide workers may continue in conservative mode, "
            "but final quality success statuses must be rejected until OCR text usability is restored."
        )
    conversion_mode = None
    worker_mode = "app-native"
    if gates:
        conversion_gate = gates.get("conversion_mode_gate") or {}
        scope_worker_gate = gates.get("scope_and_worker_gate") or {}
        conversion_mode = conversion_gate.get("user_choice")
        worker_mode = scope_worker_gate.get("worker_mode") or worker_mode
    if conversion_mode in RETIRED_CONVERSION_MODE_CHOICES or worker_mode in RETIRED_WORKER_MODE_CHOICES:
        raise ValueError(
            "Sequential token-saving execution is retired and not supported in this release. "
            "Use multi-agent-high-quality for multi-page PPTX conversion."
        )

    manifest = {
        "schema_version": "ppt-to-editable-deck-manifest-v1",
        "run_id": run_id,
        "source_pptx": str(source_pptx),
        "source_pptx_sha256": sha256_file(source_pptx),
        "source_slide_count": source_slide_count if source_slide_count is not None else len(slides),
        "slide_count": len(slides),
        "selected_slide_indices": [slide["slide_index"] for slide in slides],
        "scope": {
            "mode": "all" if (source_slide_count is None or source_slide_count == len(slides)) else "selected",
            "selected_slide_indices": [slide["slide_index"] for slide in slides],
        },
        "input_type": "image-only-pptx" if all_image_only else "unsupported-pptx",
        "slide_size": slide_size,
        "preflight": {
            "status": "passed" if all_image_only else "failed",
            "pptx_is_readable": True,
            "all_slides_are_full_page_images": all_image_only,
            "all_full_slide_images_are_png": all_png,
            "ocr_dependency_status": ocr_status["status"],
            "ocr_runtime_status": ocr_runtime_status,
            "ocr_dependency_checks": ocr_status["checks"],
            "ocr_runtime_report": ocr_status,
            "ocr_runtime": ocr_status.get("runtime"),
            "notes": notes,
        },
        "user_confirmation_gates": {
            "status": "passed" if gates else "not-recorded",
            "schema_version": gates.get("schema_version") if gates else None,
            "gates_file": str(gates_file) if gates_file else None,
            "source_pptx_sha256": gates.get("source_pptx_sha256") if gates else None,
            "conversion_mode_gate": gates.get("conversion_mode_gate") if gates else None,
            "ocr_runtime_gate": gates.get("ocr_runtime_gate") if gates else None,
            "scope_and_worker_gate": gates.get("scope_and_worker_gate") if gates else None,
        },
        "output": {
            "run_dir": str(run_dir),
            "final_deck": str(run_dir / "output" / "final-deck-fallback.pptx"),
            "deck_qa_report": str(run_dir / "output" / "deck-level-qa-report.json"),
        },
        "policy": {
            "max_attempts_per_slide": 2,
            "failed_slide_mode": "image-fallback",
            "ocr_runtime_policy": (
                "Use this package's setup_ocr_runtime/check_ocr_runtime flow. "
                "Do not rely on a local development .venv-ocr."
            ),
            "final_release_ocr_runtime_requirements": [
                "setup_ocr_runtime",
                "check_ocr_runtime",
            ],
            "gates_schema_version": GATES_SCHEMA_VERSION,
            "probe_schema_version": PROBE_SCHEMA_VERSION,
            "fallback_sticker_text": DEFAULT_STICKER_TEXT,
            "fallback_sticker_style": {
                "position": "top-right",
                "fill": "yellow",
                "border": "black",
                "text_color": "black",
            },
            "multi_agent_concurrency_cap": None,
            "multi_agent_required": True,
            "sequential_token_saving_mode": False,
            "allow_sequential_worker_identity_reuse": False,
            "worker_identity_required": True,
            "progress_heartbeat_required_statuses": REQUIRED_PROGRESS_STATUSES,
            "quality_finalize_requires_ocr_dependency_passed": True,
            "quality_finalize_requires_ocr_text_usable": True,
            "blocking_success_limitation_patterns": BLOCKING_SUCCESS_LIMITATION_PATTERNS,
            "accepted_with_limitations_policy": (
                "Only minor fidelity limits can be accepted. OCR unavailable, manual layout guess, "
                "major layout drift, rendered text overflow, font fallback, unwanted effects, or obvious text-size drift "
                "must become failed-retryable or image-fallback."
            ),
        },
        "slides": [],
    }

    for slide in slides:
        job_manifest = make_slide_job_manifest(
            run_id,
            slide,
            run_dir,
            ocr_python,
            ocr_deps,
            ocr_runtime_status=ocr_runtime_status,
            ocr_runtime_report=ocr_status,
            worker_mode=worker_mode,
            conversion_mode=conversion_mode or "multi-agent-high-quality",
        )
        picture = slide["pictures"][0] if slide["pictures"] else {}
        manifest["slides"].append(
            {
                "slide_index": slide["slide_index"],
                "source_image": picture.get("extracted_path"),
                "job_manifest": str(run_dir / "slide_jobs" / f"slide_{slide['slide_index']:02d}" / "slide_job_manifest.json"),
                "status": "pending" if slide["is_single_full_slide_picture"] else "invalid-input",
                "attempts": 0,
                "final_slide_mode": None,
                "output_pptx": None,
                "qa_summary": None,
                "failure_reason": None,
                "fallback": None,
                "limitations": [],
                "job_required_outputs": job_manifest["required_outputs"],
                "required_reading": job_manifest["required_reading"],
                "qa_summary_required_fields": job_manifest["qa_summary_required_fields"],
                "worker_contract": job_manifest["worker_contract"],
                "quality_rules": job_manifest["quality_rules"],
                "ocr_dependency_status_at_dispatch": ocr_status["status"],
                "ocr_runtime_status_at_dispatch": ocr_runtime_status,
                "required_progress_statuses": REQUIRED_PROGRESS_STATUSES,
            }
        )

    return manifest


def add_failure_sticker(slide, prs: Presentation, text: str) -> None:
    margin = Emu(int(prs.slide_width * 0.025))
    width = Emu(int(prs.slide_width * 0.18))
    height = Emu(int(prs.slide_height * 0.07))
    left = Emu(int(prs.slide_width - width - margin))
    top = margin

    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 221, 66)
    box.line.color.rgb = RGBColor(0, 0, 0)
    box.line.width = Pt(1.25)

    tf = box.text_frame
    tf.clear()
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)


def create_fallback_slide_deck(
    source_image: str,
    out_pptx: Path,
    slide_width_emu: int,
    slide_height_emu: int,
    sticker_text: str,
) -> None:
    prs = Presentation()
    prs.slide_width = Emu(slide_width_emu)
    prs.slide_height = Emu(slide_height_emu)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(source_image, 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_failure_sticker(slide, prs, sticker_text)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)


def rels_path_for_slide(slide_number: int) -> str:
    return f"ppt/slides/_rels/slide{slide_number}.xml.rels"


def next_media_name(existing_names: set[str], suffix: str, slide_number: int, counter: int) -> str:
    suffix = suffix if suffix.startswith(".") else f".{suffix}"
    while True:
        candidate = f"ppt/media/merge_slide{slide_number:02d}_{counter:03d}{suffix}"
        if candidate not in existing_names:
            existing_names.add(candidate)
            return candidate
        counter += 1


def find_next_slide_rel_id(rels_root: ET.Element) -> int:
    max_id = 0
    for rel in rels_root.findall("rel:Relationship", NS):
        rid = rel.attrib.get("Id", "")
        match = re.match(r"rId(\d+)$", rid)
        if match:
            max_id = max(max_id, int(match.group(1)))
    return max_id + 1


def find_next_slide_id(presentation_root: ET.Element) -> int:
    max_id = 255
    for sld_id in presentation_root.findall(".//p:sldId", NS):
        raw = sld_id.attrib.get("id")
        if raw and raw.isdigit():
            max_id = max(max_id, int(raw))
    return max_id + 1


def ensure_slide_override(types_root: ET.Element, slide_number: int) -> None:
    part_name = f"/ppt/slides/slide{slide_number}.xml"
    for override in types_root.findall("{http://schemas.openxmlformats.org/package/2006/content-types}Override"):
        if override.attrib.get("PartName") == part_name:
            return
    ET.SubElement(
        types_root,
        "{http://schemas.openxmlformats.org/package/2006/content-types}Override",
        {
            "PartName": part_name,
            "ContentType": "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
        },
    )


def ensure_default_content_type(types_root: ET.Element, extension: str, content_type: str) -> None:
    extension = extension.lstrip(".")
    for default in types_root.findall("{http://schemas.openxmlformats.org/package/2006/content-types}Default"):
        if default.attrib.get("Extension", "").lower() == extension.lower():
            return
    ET.SubElement(
        types_root,
        "{http://schemas.openxmlformats.org/package/2006/content-types}Default",
        {
            "Extension": extension,
            "ContentType": content_type,
        },
    )


def content_type_for_media(part_name: str) -> str | None:
    suffix = Path(part_name).suffix.lower()
    return {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".svg": "image/svg+xml",
    }.get(suffix)


def rewrite_slide_rels_and_collect_parts(
    source_zip: zipfile.ZipFile,
    source_slide_rels_name: str,
    dest_slide_number: int,
    existing_names: set[str],
) -> tuple[bytes | None, dict[str, bytes]]:
    if source_slide_rels_name not in source_zip.namelist():
        return None, {}

    rels_root = ET.fromstring(source_zip.read(source_slide_rels_name))
    extra_parts: dict[str, bytes] = {}
    media_counter = 1

    for rel in rels_root.findall("rel:Relationship", NS):
        target = rel.attrib.get("Target")
        if not target or target.startswith("http://") or target.startswith("https://"):
            continue
        source_part = resolve_rel_target(f"ppt/slides/slide1.xml", target)
        if source_part not in source_zip.namelist():
            continue

        if source_part.startswith("ppt/media/"):
            suffix = Path(source_part).suffix or ".bin"
            dest_part = next_media_name(existing_names, suffix, dest_slide_number, media_counter)
            media_counter += 1
            extra_parts[dest_part] = source_zip.read(source_part)
            rel.attrib["Target"] = f"../media/{Path(dest_part).name}"
        else:
            # Keep common layout/theme references when the base template already has them.
            # If a future single-slide packager emits custom charts/embeddings, this block
            # should be expanded into a full relationship-part copier.
            if source_part not in existing_names:
                extra_parts[source_part] = source_zip.read(source_part)
                existing_names.add(source_part)

    return ET.tostring(rels_root, encoding="utf-8", xml_declaration=True), extra_parts


def merge_single_slide_decks(slide_decks: list[Path], final_pptx: Path) -> None:
    if not slide_decks:
        raise ValueError("No slide decks to merge")

    final_pptx.parent.mkdir(parents=True, exist_ok=True)
    first = slide_decks[0]
    package_parts: dict[str, bytes] = {}
    with zipfile.ZipFile(first) as base_zip:
        for name in base_zip.namelist():
            package_parts[name] = base_zip.read(name)

    existing_names = set(package_parts)
    presentation_root = ET.fromstring(package_parts["ppt/presentation.xml"])
    rels_root = ET.fromstring(package_parts["ppt/_rels/presentation.xml.rels"])
    types_root = ET.fromstring(package_parts["[Content_Types].xml"])

    sld_id_lst = presentation_root.find("p:sldIdLst", NS)
    if sld_id_lst is None:
        raise ValueError("presentation slide id list missing")

    for slide_number, source_deck in enumerate(slide_decks[1:], start=2):
        with zipfile.ZipFile(source_deck) as source_zip:
            source_slide_name = "ppt/slides/slide1.xml"
            if source_slide_name not in source_zip.namelist():
                raise ValueError(f"{source_deck} does not contain ppt/slides/slide1.xml")

            dest_slide_name = f"ppt/slides/slide{slide_number}.xml"
            dest_rels_name = rels_path_for_slide(slide_number)
            package_parts[dest_slide_name] = source_zip.read(source_slide_name)
            existing_names.add(dest_slide_name)

            rels_bytes, extra_parts = rewrite_slide_rels_and_collect_parts(
                source_zip,
                "ppt/slides/_rels/slide1.xml.rels",
                slide_number,
                existing_names,
            )
            if rels_bytes is not None:
                package_parts[dest_rels_name] = rels_bytes
                existing_names.add(dest_rels_name)
            for part_name, part_bytes in extra_parts.items():
                package_parts.setdefault(part_name, part_bytes)
                media_content_type = content_type_for_media(part_name)
                if media_content_type:
                    ensure_default_content_type(types_root, Path(part_name).suffix, media_content_type)

            rel_id = f"rId{find_next_slide_rel_id(rels_root)}"
            ET.SubElement(
                rels_root,
                f"{{{NS['rel']}}}Relationship",
                {
                    "Id": rel_id,
                    "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
                    "Target": f"slides/slide{slide_number}.xml",
                },
            )
            ET.SubElement(
                sld_id_lst,
                f"{{{NS['p']}}}sldId",
                {
                    "id": str(find_next_slide_id(presentation_root)),
                    f"{{{NS['r']}}}id": rel_id,
                },
            )
            ensure_slide_override(types_root, slide_number)

    package_parts["ppt/presentation.xml"] = ET.tostring(
        presentation_root,
        encoding="utf-8",
        xml_declaration=True,
    )
    package_parts["ppt/_rels/presentation.xml.rels"] = ET.tostring(
        rels_root,
        encoding="utf-8",
        xml_declaration=True,
    )
    package_parts["[Content_Types].xml"] = ET.tostring(
        types_root,
        encoding="utf-8",
        xml_declaration=True,
    )

    with zipfile.ZipFile(final_pptx, "w", compression=zipfile.ZIP_DEFLATED) as out_zip:
        for name, data in package_parts.items():
            out_zip.writestr(name, data)


def mark_result_failed(result: dict, reason: str, note: str) -> None:
    result["status"] = "failed-retryable"
    result["failure_reason"] = reason
    limitations = list(result.get("limitations") or [])
    if note not in limitations:
        limitations.append(note)
    result["limitations"] = limitations


def read_progress_statuses(progress_path: Path) -> tuple[list[str], str | None]:
    if not progress_path.exists():
        return [], "progress_json_missing"
    try:
        data = read_json(progress_path)
    except json.JSONDecodeError as exc:
        return [], f"progress_json_invalid: {exc}"

    records: list[dict] = []
    if isinstance(data, list):
        records = [item for item in data if isinstance(item, dict)]
    elif isinstance(data, dict):
        for key in ("history", "events", "milestones", "checkpoints", "progress"):
            value = data.get(key)
            if isinstance(value, list):
                records = [item for item in value if isinstance(item, dict)]
                break
        if not records and isinstance(data.get("status"), str):
            records = [data]
    else:
        return [], "progress_json_must_be_object_or_list"

    statuses = [str(item.get("status", "")) for item in records if item.get("status")]
    return statuses, None


def validate_progress_heartbeat(progress_path: Path, required_statuses: list[str]) -> tuple[bool, str | None]:
    statuses, error = read_progress_statuses(progress_path)
    if error:
        return False, error
    missing = [status for status in required_statuses if status not in statuses]
    if missing:
        return False, "progress_heartbeat_missing: " + ", ".join(missing)
    return True, None


def validate_worker_identity(
    summary: dict,
    slide_index: int,
    allowed_execution_models: list[str] | None = None,
) -> tuple[bool, str | None, str | None]:
    raw = summary.get("worker_identity")
    if raw is None:
        # Backward-compatible aliases are accepted, but new workers should use worker_identity.
        aliases = {
            "worker_thread_id": summary.get("worker_thread_id"),
            "worker_agent_id": summary.get("worker_agent_id"),
            "execution_model": summary.get("execution_model"),
            "slide_scope": summary.get("slide_scope"),
        }
        raw = {k: v for k, v in aliases.items() if v is not None}
    if not isinstance(raw, dict) or not raw:
        return False, "worker_identity_missing", None

    worker_id = raw.get("worker_thread_id") or raw.get("worker_agent_id")
    if not worker_id:
        return False, "worker_identity_missing_worker_id", None

    execution_model = str(raw.get("execution_model", ""))
    allowed_models = allowed_execution_models or ["delegated-real-per-slide-worker"]
    if execution_model and execution_model not in allowed_models:
        return False, "worker_identity_wrong_execution_model", str(worker_id)

    scope = raw.get("slide_scope", raw.get("processed_slide_indices", raw.get("source_slide_index")))
    if scope is None:
        return False, "worker_identity_missing_slide_scope", str(worker_id)
    if isinstance(scope, int):
        scoped = [scope]
    elif isinstance(scope, str) and scope.isdigit():
        scoped = [int(scope)]
    elif isinstance(scope, list):
        scoped = []
        for item in scope:
            if isinstance(item, int):
                scoped.append(item)
            elif isinstance(item, str) and item.isdigit():
                scoped.append(int(item))
    else:
        return False, "worker_identity_invalid_slide_scope", str(worker_id)
    if scoped != [slide_index]:
        return False, "worker_identity_scope_not_single_current_slide", str(worker_id)

    return True, None, str(worker_id)


def flatten_quality_text(value) -> list[str]:
    if value is None:
        return []
    if isinstance(value, str):
        return [value]
    if isinstance(value, list):
        texts: list[str] = []
        for item in value:
            texts.extend(flatten_quality_text(item))
        return texts
    if isinstance(value, dict):
        texts = []
        for item in value.values():
            texts.extend(flatten_quality_text(item))
        return texts
    return [str(value)]


def find_blocking_success_limitations(summary: dict) -> list[str]:
    candidates: list[str] = []
    candidates.extend(flatten_quality_text(summary.get("limitations")))
    visual_qa = summary.get("visual_qa")
    if isinstance(visual_qa, dict):
        candidates.extend(flatten_quality_text(visual_qa.get("known_visual_risks")))
        candidates.extend(flatten_quality_text(visual_qa.get("remaining_visual_issues")))
        candidates.extend(flatten_quality_text(visual_qa.get("failed_checks")))
    matches: list[str] = []
    for text in candidates:
        normalized = text.lower()
        for pattern in BLOCKING_SUCCESS_LIMITATION_PATTERNS:
            if pattern in normalized:
                matches.append(text)
                break
    return matches


def script_name_from_reported_value(value) -> str | None:
    if isinstance(value, dict):
        for key in ("name", "script", "path", "script_path", "command"):
            if key in value:
                return script_name_from_reported_value(value[key])
        return None
    if value is None:
        return None
    text = str(value).strip()
    if not text:
        return None
    return text.replace("\\", "/").split("/")[-1]


def validate_single_page_engine_usage(summary: dict) -> tuple[bool, str | None]:
    engine_used = str(summary.get("single_page_engine_used", "")).lower()
    if PRIMARY_SINGLE_PAGE_ENGINE_ID not in engine_used and "ppt-to-editable-3-0" not in engine_used:
        return False, "single_page_engine_used_missing_or_not_3_0"

    raw_scripts = summary.get("single_page_engine_scripts_used")
    if not isinstance(raw_scripts, list) or not raw_scripts:
        return False, "single_page_engine_scripts_used_missing"

    script_names = {
        name
        for item in raw_scripts
        if (name := script_name_from_reported_value(item))
    }
    if not script_names:
        return False, "single_page_engine_scripts_used_empty"

    mode = str(summary.get("mode", "")).lower()
    if "reconstruction" in mode:
        required = RECONSTRUCTION_PACKAGING_SCRIPTS
        if not (script_names & required):
            return False, "reconstruction_without_integrated_packager_or_qa_script"
    elif "clean-background" in mode or "clean_background" in mode:
        required = CLEAN_BACKGROUND_PACKAGING_SCRIPTS
        if not (script_names & required):
            return False, "clean_background_without_integrated_packager_script"
    elif not (script_names & (RECONSTRUCTION_PACKAGING_SCRIPTS | CLEAN_BACKGROUND_PACKAGING_SCRIPTS)):
        return False, "success_without_integrated_packaging_script"

    return True, None


def resolve_attempt_artifact_path(attempt_dir: Path, value: str | None, default_name: str) -> Path:
    if value:
        candidate = Path(value)
        if not candidate.is_absolute():
            candidate = attempt_dir / candidate
        return candidate
    return attempt_dir / default_name


def validate_render_report(attempt_dir: Path, summary: dict) -> tuple[bool, str | None]:
    report_value = summary.get("render_report")
    if not report_value:
        visual_qa = summary.get("visual_qa")
        if isinstance(visual_qa, dict):
            report_value = visual_qa.get("render_report")
    report_path = resolve_attempt_artifact_path(
        attempt_dir,
        str(report_value) if report_value else None,
        "render_report.json",
    )
    if not report_path.exists():
        return False, f"render_report_missing: {report_path}"
    try:
        report = read_json(report_path)
    except json.JSONDecodeError as exc:
        return False, f"render_report_invalid_json: {exc}"
    if report.get("status") != "rendered":
        return False, f"render_report_not_rendered: {report.get('status')}"
    rendered_png = report.get("rendered_png")
    if not rendered_png:
        return False, "render_report_missing_rendered_png"
    rendered_path = Path(rendered_png)
    if not rendered_path.is_absolute():
        rendered_path = report_path.parent / rendered_path
    if not rendered_path.exists():
        return False, f"rendered_png_missing: {rendered_path}"
    return True, None


def review_has_blocking_issue(review: dict, patterns: tuple[str, ...]) -> str | None:
    candidates: list[str] = []
    candidates.extend(flatten_quality_text(review.get("issues")))
    candidates.extend(flatten_quality_text(review.get("warnings")))
    candidates.extend(flatten_quality_text(review.get("failed_checks")))
    candidates.extend(flatten_quality_text(review.get("remaining_issues")))
    candidates.extend(flatten_quality_text(review.get("status")))
    for text in candidates:
        normalized = text.lower()
        for pattern in patterns:
            if pattern in normalized:
                return text
    return None


def validate_font_size_consistency_review(summary: dict) -> tuple[bool, str | None]:
    review = summary.get("font_size_consistency_review")
    if not isinstance(review, dict):
        return False, "font_size_consistency_review_missing"
    status = str(review.get("status", "")).lower()
    if status in {"failed", "fail", "not-reviewed", "missing"}:
        return False, f"font_size_consistency_review_{status or 'invalid'}"
    blocking = review_has_blocking_issue(
        review,
        ("tiny-font-outlier", "font-size-hierarchy-drift", "unnecessary-font-shrink"),
    )
    if blocking:
        return False, "font_size_consistency_blocking_issue: " + blocking
    return True, None


def validate_complex_visual_strategy_review(summary: dict) -> tuple[bool, str | None]:
    review = summary.get("complex_visual_strategy_review")
    if not isinstance(review, dict):
        return False, "complex_visual_strategy_review_missing"
    status = str(review.get("status", "")).lower()
    if status in {"failed", "fail", "not-reviewed", "missing"}:
        return False, f"complex_visual_strategy_review_{status or 'invalid'}"
    if review.get("unresolved_native_redraw_drift") is True:
        return False, "complex_visual_strategy_unresolved_native_redraw_drift"
    blocking = review_has_blocking_issue(
        review,
        ("native-redraw-complex-artwork", "complex artwork native redraw drift", "native-redraw-worsened-fidelity"),
    )
    if blocking:
        return False, "complex_visual_strategy_blocking_issue: " + blocking
    return True, None


def validate_text_background_policy(summary: dict) -> tuple[bool, str | None]:
    policy = summary.get("text_background_policy")
    if not isinstance(policy, dict):
        return False, "text_background_policy_missing"
    status = str(policy.get("status", "")).lower()
    if status in {"failed", "fail", "not-reviewed", "missing"}:
        return False, f"text_background_policy_{status or 'invalid'}"
    if policy.get("invented_text_backgrounds_added") is True:
        return False, "invented_text_backgrounds_added"
    blocking = review_has_blocking_issue(
        policy,
        ("invented-text-background", "pale rectangle", "text backing rectangle", "backing rectangles create"),
    )
    if blocking:
        return False, "text_background_policy_blocking_issue: " + blocking
    return True, None


def validate_non_text_anchor_preservation_review(summary: dict) -> tuple[bool, str | None]:
    review = summary.get("non_text_anchor_preservation_review")
    if not isinstance(review, dict):
        return False, "non_text_anchor_preservation_review_missing"
    status = str(review.get("status", "")).lower()
    if status in {"failed", "fail", "not-reviewed", "missing"}:
        return False, f"non_text_anchor_preservation_review_{status or 'invalid'}"
    if review.get("missing_required_anchors") is True:
        return False, "non_text_anchor_missing_required_anchors"
    blocking = review_has_blocking_issue(
        review,
        (
            "missing-non-text-anchor",
            "missing-arrow",
            "missing-arrowhead",
            "erased-arrow",
            "broken-route-line",
            "broken-connector",
            "missing-connector",
            "missing-step-marker",
            "missing-card-edge",
            "cropped-icon",
            "icon stroke missing",
        ),
    )
    if blocking:
        return False, "non_text_anchor_preservation_blocking_issue: " + blocking
    return True, None


def is_raw_single_page_summary(summary: dict, slide_info: dict) -> bool:
    if not isinstance(summary, dict):
        return False
    if summary.get("schema_version") == "ppt-to-editable-3-0-slide-qa-summary-v1":
        return False
    raw_status = str(summary.get("status", "")).lower()
    raw_status_like = raw_status in {"pass", "passed", "success", "warning", "fail", "failed", "generated"}
    has_raw_single_page_markers = any(key in summary for key in ("steps", "artifacts", "checks", "summary"))
    required_fields = slide_info.get("qa_summary_required_fields", [])
    missing_fields = [field for field in required_fields if field not in summary]
    return raw_status_like and (has_raw_single_page_markers or bool(missing_fields))


def autoclose_single_page_summary_if_needed(
    attempt_dir: Path,
    slide_info: dict,
    run_dir: Path,
    summary: dict,
) -> tuple[dict, str | None]:
    if not is_raw_single_page_summary(summary, slide_info):
        return summary, None

    slide_index = slide_info["slide_index"]
    handoff_script = Path(__file__).with_name("write_slide_qa_summary_from_single_page.py")
    manifest_path = run_dir / "slide_jobs" / f"slide_{slide_index:02d}" / "slide_job_manifest.json"
    if not handoff_script.exists():
        return summary, f"single_page_handoff_script_missing: {handoff_script}"
    if not manifest_path.exists():
        return summary, f"single_page_handoff_manifest_missing: {manifest_path}"

    worker_identity = summary.get("worker_identity") if isinstance(summary.get("worker_identity"), dict) else {}
    worker_id = (
        worker_identity.get("worker_thread_id")
        or worker_identity.get("worker_agent_id")
        or summary.get("worker_thread_id")
        or summary.get("worker_agent_id")
        or f"controller-autoclosed-slide-{slide_index:02d}"
    )
    proc = subprocess.run(
        [
            sys.executable,
            str(handoff_script),
            "--slide-job-manifest",
            str(manifest_path),
            "--attempt-dir",
            str(attempt_dir),
            "--worker-agent-id",
            str(worker_id),
            "--output",
            str(attempt_dir / "qa_summary.json"),
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
    )
    if proc.returncode != 0:
        details = (proc.stderr or proc.stdout or "").strip()
        return summary, f"single_page_handoff_autoclose_failed: {details or proc.returncode}"

    try:
        return read_json(attempt_dir / "qa_summary.json"), None
    except json.JSONDecodeError as exc:
        return summary, f"single_page_handoff_autoclose_invalid_json: {exc}"


def collect_slide_result(slide_info: dict, run_dir: Path) -> dict:
    slide_index = slide_info["slide_index"]
    slide_dir = run_dir / "slide_jobs" / f"slide_{slide_index:02d}"
    attempts = sorted(slide_dir.glob("attempt_*"))
    latest = attempts[-1] if attempts else slide_dir / "attempt_01"
    qa_summary = latest / "qa_summary.json"
    slide_pptx = latest / "slide.pptx"
    progress_path = latest / "progress.json"

    result = {
        "attempt_dir": str(latest),
        "qa_summary_path": str(qa_summary) if qa_summary.exists() else None,
        "slide_pptx": str(slide_pptx) if slide_pptx.exists() else None,
        "progress_path": str(progress_path) if progress_path.exists() else None,
        "status": "missing-output",
        "failure_reason": "single_slide_agent_output_missing",
        "limitations": [],
        "qa_summary": None,
        "worker_key": None,
    }

    if qa_summary.exists():
        try:
            summary = read_json(qa_summary)
        except json.JSONDecodeError as exc:
            result["failure_reason"] = f"qa_summary_invalid_json: {exc}"
            return result
        summary, autoclose_error = autoclose_single_page_summary_if_needed(latest, slide_info, run_dir, summary)
        if autoclose_error:
            result["failure_reason"] = autoclose_error
            return result
        result["qa_summary"] = summary
        status = summary.get("status", "missing-status")
        result["status"] = status
        result["failure_reason"] = summary.get("failure_reason")
        result["limitations"] = summary.get("limitations", [])

        required_fields = slide_info.get("qa_summary_required_fields", [])
        missing_fields = [field for field in required_fields if field not in summary]
        if missing_fields:
            result["status"] = "failed-retryable"
            result["failure_reason"] = "missing_required_qa_summary_fields: " + ", ".join(missing_fields)
            limitations = list(result["limitations"])
            limitations.append(
                    "Controller rejected the slide output because qa_summary.json did not report all required v3.1 preview fields."
            )
            result["limitations"] = limitations

        required_ids = {
            item["id"]
            for item in slide_info.get("required_reading", [])
            if item.get("must_read")
        }
        raw_read = summary.get("references_read", [])
        if isinstance(raw_read, list):
            read_ids = {
                item.get("id") if isinstance(item, dict) else str(item)
                for item in raw_read
            }
        else:
            read_ids = set()
        missing_ids = sorted(required_ids - read_ids)
        if missing_ids:
            result["status"] = "failed-retryable"
            result["failure_reason"] = "missing_required_references_read: " + ", ".join(missing_ids)
            limitations = list(result["limitations"])
            limitations.append("Controller rejected the slide output because required references were not reported in qa_summary.references_read.")
            result["limitations"] = limitations

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            engine_ok, engine_error = validate_single_page_engine_usage(summary)
            if not engine_ok:
                mark_result_failed(
                    result,
                    engine_error or "single_page_engine_usage_invalid",
                    "Controller rejected the slide output because it did not prove use of the bundled v3.1 preview single-page engine and deterministic packaging scripts.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            worker_contract = slide_info.get("worker_contract") if isinstance(slide_info.get("worker_contract"), dict) else {}
            allowed_execution_models = worker_contract.get("allowed_execution_models")
            worker_ok, worker_error, worker_key = validate_worker_identity(
                summary,
                slide_index,
                allowed_execution_models if isinstance(allowed_execution_models, list) else None,
            )
            result["worker_key"] = worker_key
            if not worker_ok:
                mark_result_failed(
                    result,
                    worker_error or "worker_identity_invalid",
                    "Controller rejected the slide output because it did not prove a real separate per-slide worker.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            progress_ok, progress_error = validate_progress_heartbeat(
                progress_path,
                slide_info.get("required_progress_statuses", REQUIRED_PROGRESS_STATUSES),
            )
            if not progress_ok:
                mark_result_failed(
                    result,
                    progress_error or "progress_heartbeat_invalid",
                    "Controller rejected the slide output because progress.json did not contain all required heartbeat milestones.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            if slide_info.get("ocr_dependency_status_at_dispatch") != "passed":
                mark_result_failed(
                    result,
                    "ocr_runtime_unavailable_for_quality_finalize",
                    "Controller rejected the slide output because OCR runtime was unavailable at dispatch; quality runs cannot accept manual-layout reconstruction as success.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            dispatch_ocr_status = normalize_ocr_runtime_status(
                slide_info.get("ocr_runtime_status_at_dispatch"),
                slide_info.get("ocr_dependency_status_at_dispatch"),
            )
            summary_ocr_status = normalize_ocr_runtime_status(
                summary.get("ocr_runtime_status"),
                slide_info.get("ocr_dependency_status_at_dispatch"),
            )
            if not ocr_text_is_usable(dispatch_ocr_status) or not ocr_text_is_usable(summary_ocr_status):
                mark_result_failed(
                    result,
                    (
                        "ocr_text_not_usable_for_success: "
                        f"dispatch={dispatch_ocr_status}; qa_summary={summary_ocr_status}"
                    ),
                    "Controller rejected the slide output because OCR Chinese text was not proven usable; geometry-only OCR cannot support a success status.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            blocking = find_blocking_success_limitations(summary)
            if blocking:
                mark_result_failed(
                    result,
                    "blocking_limitations_not_allowed_for_success: " + "; ".join(blocking[:5]),
                    "Controller rejected the slide output because success statuses cannot include OCR-unavailable/manual-layout/major-drift limitations.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            if summary.get("rendered_png_reviewed") is not True:
                mark_result_failed(
                    result,
                    "success_without_rendered_png_review",
                    "Controller rejected success status because rendered PowerPoint/LibreOffice PNG review was not reported.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            render_ok, render_error = validate_render_report(latest, summary)
            if not render_ok:
                mark_result_failed(
                    result,
                    render_error or "render_report_invalid",
                    "Controller rejected success status because standard render_report.json did not prove a rendered PNG exists.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            region_plan = summary.get("region_crop_plan")
            if not isinstance(region_plan, list) or len(region_plan) < 3:
                mark_result_failed(
                    result,
                    "region_crop_plan_missing_or_too_shallow",
                    "Controller rejected the slide output because complex structured pages need a semantic 3-5 region crop/native/text plan.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            region_plan_ok, region_plan_error = validate_region_crop_plan_contract(summary)
            if not region_plan_ok:
                mark_result_failed(
                    result,
                    region_plan_error or "region_crop_plan_contract_failed",
                    "Controller rejected the slide output because region_crop_plan must not contain not-reported placeholders and each region needs source bbox plus rationale.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            strategy_ok, strategy_error = validate_visual_strategy_artifacts(latest, summary)
            if not strategy_ok:
                mark_result_failed(
                    result,
                    strategy_error or "visual_strategy_artifact_gate_failed",
                    "Controller rejected the slide output because source_reconstruction_plan.json or visual_strategy_report.json failed crop-first visual strategy evidence checks.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            font_policy = str(summary.get("font_policy_applied", "")).lower()
            if "yahei" not in font_policy and "微软雅黑" not in font_policy:
                result["status"] = "failed-retryable"
                result["failure_reason"] = "font_policy_not_applied_or_not_reported"
                limitations = list(result["limitations"])
                limitations.append("Controller rejected the slide output because Microsoft YaHei font policy was not reported.")
                result["limitations"] = limitations

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations") and summary.get("effects_cleared") is not True:
            result["status"] = "failed-retryable"
            result["failure_reason"] = "effects_cleared_not_true"
            limitations = list(result["limitations"])
            limitations.append("Controller rejected the slide output because native shape effect clearing was not reported as true.")
            result["limitations"] = limitations

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            font_review_ok, font_review_error = validate_font_size_consistency_review(summary)
            if not font_review_ok:
                mark_result_failed(
                    result,
                    font_review_error or "font_size_consistency_review_invalid",
                    "Controller rejected the slide output because font size consistency was not reviewed or a tiny-font drift was reported.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            complex_review_ok, complex_review_error = validate_complex_visual_strategy_review(summary)
            if not complex_review_ok:
                mark_result_failed(
                    result,
                    complex_review_error or "complex_visual_strategy_review_invalid",
                    "Controller rejected the slide output because complex visual regions were not crop-first reviewed or native redraw drift was reported.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            background_policy_ok, background_policy_error = validate_text_background_policy(summary)
            if not background_policy_ok:
                mark_result_failed(
                    result,
                    background_policy_error or "text_background_policy_invalid",
                    "Controller rejected the slide output because invented text backing fills are not allowed.",
                )

        status_for_policy = result["status"]
        if status_for_policy in ("passed", "accepted-with-limitations"):
            anchor_review_ok, anchor_review_error = validate_non_text_anchor_preservation_review(summary)
            if not anchor_review_ok:
                mark_result_failed(
                    result,
                    anchor_review_error or "non_text_anchor_preservation_review_invalid",
                    "Controller rejected the slide output because non-text visual anchors were missing, broken, cropped, or erased.",
                )

    if result["status"] in ("passed", "accepted-with-limitations") and not slide_pptx.exists():
        result["status"] = "missing-output"
        result["failure_reason"] = "qa_summary_passed_but_slide_pptx_missing"

    return result


def assemble_final_deck_from_results(manifest: dict, run_dir: Path) -> dict:
    output_dir = run_dir / "output"
    output_dir.mkdir(parents=True, exist_ok=True)
    per_slide_dir = output_dir / "per-slide-decks"
    per_slide_dir.mkdir(parents=True, exist_ok=True)

    slide_decks: list[Path] = []
    fallback_count = 0
    passed_count = 0
    accepted_count = 0
    fallback_text = manifest["policy"]["fallback_sticker_text"]
    collected_results = []

    for slide_info in manifest["slides"]:
        result = collect_slide_result(slide_info, run_dir)
        collected_results.append((slide_info, result))

    if manifest.get("policy", {}).get("allow_sequential_worker_identity_reuse") is not True:
        successful_by_worker: dict[str, list[dict]] = {}
        for _slide_info, result in collected_results:
            if result["status"] in ("passed", "accepted-with-limitations") and result.get("worker_key"):
                successful_by_worker.setdefault(result["worker_key"], []).append(result)
        for worker_key, results in successful_by_worker.items():
            if len(results) > 1:
                for result in results:
                    mark_result_failed(
                        result,
                        f"duplicate_worker_identity_across_slides: {worker_key}",
                        "Controller rejected the slide output because the same worker identity was used for more than one slide.",
                    )

    for slide_info, result in collected_results:
        status = result["status"]
        if status in ("passed", "accepted-with-limitations") and result["slide_pptx"]:
            slide_decks.append(Path(result["slide_pptx"]))
            slide_info["status"] = status
            slide_info["attempts"] = max(slide_info["attempts"], 1)
            slide_info["final_slide_mode"] = "editable-pptx"
            slide_info["output_pptx"] = result["slide_pptx"]
            slide_info["qa_summary"] = result["qa_summary_path"]
            slide_info["failure_reason"] = None
            slide_info["limitations"] = result["limitations"]
            if status == "passed":
                passed_count += 1
            else:
                accepted_count += 1
        else:
            fallback_pptx = per_slide_dir / f"slide_{slide_info['slide_index']:02d}_fallback.pptx"
            create_fallback_slide_deck(
                slide_info["source_image"],
                fallback_pptx,
                manifest["slide_size"]["width_emu"],
                manifest["slide_size"]["height_emu"],
                fallback_text,
            )
            slide_decks.append(fallback_pptx)
            fallback_count += 1
            slide_info["status"] = "image-fallback"
            slide_info["attempts"] = max(slide_info["attempts"], 2)
            slide_info["final_slide_mode"] = "image-fallback"
            slide_info["output_pptx"] = str(fallback_pptx)
            slide_info["qa_summary"] = result["qa_summary_path"]
            slide_info["failure_reason"] = result["failure_reason"]
            slide_info["fallback"] = {
                "source_image": slide_info["source_image"],
                "sticker_text": fallback_text,
                "sticker_position": "top-right",
            }

    final_deck = output_dir / "final-deck.pptx"
    merge_single_slide_decks(slide_decks, final_deck)
    manifest["output"]["final_deck"] = str(final_deck)

    report = {
        "schema_version": "ppt-to-editable-deck-qa-report-v1",
        "run_id": manifest["run_id"],
        "source_pptx": manifest["source_pptx"],
        "final_deck": str(final_deck.resolve()),
        "slide_count": manifest["slide_count"],
        "passed_count": passed_count,
        "accepted_with_limitations_count": accepted_count,
        "image_fallback_count": fallback_count,
        "fallback_sticker_editable_text_bodies_count": fallback_count,
        "structural_editability_note": (
            "Editable text bodies on fallback pages are failure stickers, "
            "not successful slide-content conversion."
        ),
        "page_order_preserved": True,
        "fallback_sticker": manifest["policy"]["fallback_sticker_style"]
        | {"text": manifest["policy"]["fallback_sticker_text"]},
        "slides": [
            {
                "slide_index": s["slide_index"],
                "status": s["status"],
                "final_slide_mode": s["final_slide_mode"],
                "source_image": s["source_image"],
                "output_pptx": s["output_pptx"],
                "failure_reason": s["failure_reason"],
                "limitations": s["limitations"],
            }
            for s in manifest["slides"]
        ],
    }
    report_path = output_dir / "deck-level-qa-report.json"
    write_json(report_path, report)
    manifest["output"]["deck_qa_report"] = str(report_path)
    return report


def assemble_fallback_deck(manifest: dict, run_dir: Path) -> dict:
    prs = Presentation()
    prs.slide_width = Emu(manifest["slide_size"]["width_emu"])
    prs.slide_height = Emu(manifest["slide_size"]["height_emu"])
    blank = prs.slide_layouts[6]

    fallback_text = manifest["policy"]["fallback_sticker_text"]
    fallback_count = 0
    for slide_info in manifest["slides"]:
        slide = prs.slides.add_slide(blank)
        source_image = slide_info["source_image"]
        slide.shapes.add_picture(source_image, 0, 0, width=prs.slide_width, height=prs.slide_height)

        if slide_info["status"] not in ("passed", "accepted-with-limitations"):
            add_failure_sticker(slide, prs, fallback_text)
            fallback_count += 1
            slide_info["status"] = "image-fallback"
            slide_info["attempts"] = max(slide_info["attempts"], 2)
            slide_info["final_slide_mode"] = "image-fallback"
            slide_info["failure_reason"] = "no_single_slide_editable_output_available"
            slide_info["fallback"] = {
                "source_image": source_image,
                "sticker_text": fallback_text,
                "sticker_position": "top-right",
            }

    output_dir = run_dir / "output"
    output_dir.mkdir(parents=True, exist_ok=True)
    final_deck = output_dir / "final-deck-fallback.pptx"
    prs.save(final_deck)
    manifest["output"]["final_deck"] = str(final_deck)

    report = {
        "schema_version": "ppt-to-editable-deck-qa-report-v1",
        "run_id": manifest["run_id"],
        "source_pptx": manifest["source_pptx"],
        "final_deck": str(final_deck.resolve()),
        "slide_count": manifest["slide_count"],
        "passed_count": sum(1 for s in manifest["slides"] if s["status"] == "passed"),
        "accepted_with_limitations_count": sum(
            1 for s in manifest["slides"] if s["status"] == "accepted-with-limitations"
        ),
        "image_fallback_count": fallback_count,
        "fallback_sticker_editable_text_bodies_count": fallback_count,
        "structural_editability_note": (
            "Editable text bodies in an all-fallback deck are failure stickers, "
            "not successful slide-content conversion."
        ),
        "page_order_preserved": True,
        "fallback_sticker": manifest["policy"]["fallback_sticker_style"]
        | {"text": manifest["policy"]["fallback_sticker_text"]},
        "slides": [
            {
                "slide_index": s["slide_index"],
                "status": s["status"],
                "final_slide_mode": s["final_slide_mode"],
                "source_image": s["source_image"],
                "output_pptx": s["output_pptx"],
                "failure_reason": s["failure_reason"],
                "limitations": s["limitations"],
            }
            for s in manifest["slides"]
        ],
    }
    report_path = output_dir / "deck-level-qa-report.json"
    write_json(report_path, report)
    manifest["output"]["deck_qa_report"] = str(report_path)
    return report


def create_run(
    source_pptx: Path,
    run_dir: Path,
    finalize_fallback: bool,
    finalize: bool,
    ocr_python: Path | None = None,
    ocr_deps: Path | None = None,
    slides: str | None = None,
    all_slides: bool = False,
    gates_file: Path | None = None,
    allow_ocr_not_ready: bool = False,
) -> dict:
    source_pptx = source_pptx.resolve()
    run_dir = run_dir.resolve()
    if not source_pptx.exists():
        raise FileNotFoundError(source_pptx)
    if source_pptx.suffix.lower() != ".pptx":
        raise ValueError("Only .pptx input is supported by the deck controller")

    manifest_path = run_dir / "deck_manifest.json"
    if (finalize or finalize_fallback) and manifest_path.exists():
        manifest = read_json(manifest_path)
    else:
        gates, resolved_gates_file = load_and_validate_gates(
            gates_file,
            source_pptx,
            slides,
            all_slides,
            allow_ocr_not_ready,
        )
        run_id = run_dir.name
        selected_slide_indices = require_explicit_scope(slides, all_slides)
        run_dir.mkdir(parents=True, exist_ok=True)
        slide_size, all_slide_infos = extract_image_only_slides(source_pptx, run_dir)
        scoped_slides = filter_slides_by_selection(all_slide_infos, selected_slide_indices)
        manifest = build_deck_manifest(
            source_pptx,
            run_dir,
            run_id,
            slide_size,
            scoped_slides,
            source_slide_count=len(all_slide_infos),
            ocr_python=ocr_python,
            ocr_deps=ocr_deps,
            allow_ocr_not_ready=allow_ocr_not_ready,
            gates=gates,
            gates_file=resolved_gates_file,
        )
        write_json(manifest_path, manifest)

    report = None
    if finalize:
        report = assemble_final_deck_from_results(manifest, run_dir)
        write_json(manifest_path, manifest)
    elif finalize_fallback:
        report = assemble_fallback_deck(manifest, run_dir)
        write_json(manifest_path, manifest)

    return {
        "run_dir": str(run_dir),
        "deck_manifest": str(manifest_path),
        "source_slide_count": manifest.get("source_slide_count", manifest["slide_count"]),
        "slide_count": manifest["slide_count"],
        "selected_slide_indices": manifest.get("selected_slide_indices", []),
        "preflight_status": manifest["preflight"]["status"],
        "ocr_dependency_status": manifest["preflight"]["ocr_dependency_status"],
        "ocr_runtime_status": manifest["preflight"]["ocr_runtime_status"],
        "user_confirmation_gates_status": manifest.get("user_confirmation_gates", {}).get("status"),
        "all_slides_are_full_page_images": manifest["preflight"]["all_slides_are_full_page_images"],
        "final_deck": manifest["output"]["final_deck"] if (finalize or finalize_fallback) else None,
        "deck_qa_report": manifest["output"]["deck_qa_report"] if report else None,
    }


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="ppt-to-editable v3.1 preview deck controller")
    parser.add_argument("source_pptx", type=Path)
    parser.add_argument("--run-dir", type=Path, default=None)
    parser.add_argument("--ocr-python", type=Path, default=None)
    parser.add_argument("--ocr-deps", type=Path, default=None)
    parser.add_argument(
        "--probe",
        action="store_true",
        help="Read-only preflight probe. Prints deck/OCR status and a gates file template; creates no run directory.",
    )
    parser.add_argument(
        "--gates-file",
        type=Path,
        default=None,
        help="JSON evidence that Conversion Mode, OCR Runtime, and Scope gates were explained and confirmed by the user.",
    )
    scope = parser.add_mutually_exclusive_group()
    scope.add_argument(
        "--slides",
        default=None,
        help="Selected slide scope confirmed by the user, such as '1', '1,3,5', or '2-4'.",
    )
    scope.add_argument(
        "--all-slides",
        action="store_true",
        help="Convert all slides. Use only after the user explicitly confirms full-deck scope.",
    )
    parser.add_argument(
        "--user-accepted-ocr-limited-run",
        action="store_true",
        help=(
            "Bypass the OCR text-usability hard gate only after the user explicitly accepts "
            "a diagnostic or fallback-only run. Do not use for quality editable conversion."
        ),
    )
    parser.add_argument(
        "--finalize",
        action="store_true",
        help="Collect per-slide Agent outputs and assemble passed slides plus fallback pages.",
    )
    parser.add_argument(
        "--finalize-fallback",
        action="store_true",
        help="Create a complete final deck using original PNGs plus failure stickers for unresolved slides.",
    )
    args = parser.parse_args(argv)
    if args.probe:
        if args.finalize or args.finalize_fallback:
            parser.error("--probe cannot be combined with --finalize or --finalize-fallback")
    elif args.run_dir is None:
        parser.error("--run-dir is required unless --probe is used")
    return args


def main(argv: list[str]) -> int:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass
    args = parse_args(argv)
    if args.probe:
        result = probe_deck(args.source_pptx, args.ocr_python, args.ocr_deps)
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0

    result = create_run(
        args.source_pptx,
        args.run_dir,
        args.finalize_fallback,
        args.finalize,
        args.ocr_python,
        args.ocr_deps,
        args.slides,
        args.all_slides,
        args.gates_file,
        args.user_accepted_ocr_limited_run,
    )
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
