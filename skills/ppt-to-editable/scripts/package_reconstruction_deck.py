#!/usr/bin/env python3
"""MVP deterministic reconstruction packager for Round 1.

Reads a source_reconstruction_plan.json and creates a PPTX with native shapes,
lines, and editable text boxes. This is intentionally small: it proves the
plan-driven packaging path before broader OCR/visual fidelity work.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any


def add_pptx_deps(path: Path | None) -> None:
    if path and path.exists():
        sys.path.insert(0, str(path.resolve()))


def rgb(value: str):
    from pptx.dml.color import RGBColor

    value = str(value or "000000").strip().lstrip("#")
    if len(value) != 6:
        value = "000000"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def inches(value: float):
    from pptx.util import Inches

    return Inches(float(value))


def pt(value: float):
    from pptx.util import Pt

    return Pt(float(value))


def as_bool(value: Any, default: bool = False) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    normalized = str(value).strip().lower()
    if normalized in {"1", "true", "yes", "y", "on"}:
        return True
    if normalized in {"0", "false", "no", "n", "off"}:
        return False
    return default


def xml_local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def clear_inherited_effects(shape, counts: dict[str, int]) -> None:
    """Remove theme-inherited style/effect XML that can add shadows or glows."""
    removed = 0
    root = shape._element
    for parent in list(root.iter()):
        for child in list(parent):
            if xml_local_name(child.tag) in {"style", "effectLst", "effectDag"}:
                parent.remove(child)
                removed += 1
    if removed:
        counts["effectClearedObjects"] += 1
        counts["effectXmlNodesCleared"] += removed


def source_kind(element: dict[str, Any]) -> str:
    source = element.get("source")
    if isinstance(source, dict) and source.get("kind"):
        return str(source["kind"]).lower()
    if element.get("geometry_source"):
        return str(element["geometry_source"]).lower()
    return "manual"


def trace_level(element: dict[str, Any]) -> str:
    if element.get("trace_level"):
        return str(element["trace_level"]).lower()
    return "paragraph" if as_bool(element.get("semantic_block"), True) else "line"


def text_overflow_risk(element: dict[str, Any], lines: list[str]) -> bool:
    """Cheap fit warning; final judgment still comes from rendered QA."""
    font_size = float(element.get("font_size", 14))
    margin = float(element.get("margin", 0.03))
    box_w = max(float(element.get("w", 0)) - margin * 2, 0.01)
    box_h = max(float(element.get("h", 0)) - margin * 2, 0.01)
    needed_h = len(lines) * font_size * 1.18 / 72.0
    if needed_h > box_h:
        return True

    for line in lines:
        width_units = sum(1.0 if ord(ch) > 127 else 0.55 for ch in line)
        needed_w = width_units * font_size * 0.58 / 72.0
        if needed_w > box_w and not as_bool(element.get("word_wrap"), True):
            return True
    return False


def resolve_resource(path_value: str, plan_dir: Path) -> Path:
    path = Path(path_value)
    if path.is_absolute():
        return path
    plan_relative = plan_dir / path
    if plan_relative.exists():
        return plan_relative
    cwd_relative = Path.cwd() / path
    if cwd_relative.exists():
        return cwd_relative
    return plan_relative


def shape_type(name: str):
    from pptx.enum.shapes import MSO_SHAPE

    normalized = str(name or "rect").lower()
    if normalized in {"roundrect", "rounded_rectangle", "rounded-rectangle"}:
        return MSO_SHAPE.ROUNDED_RECTANGLE
    if normalized in {"oval", "ellipse", "circle"}:
        return MSO_SHAPE.OVAL
    if normalized in {"chevron"}:
        return MSO_SHAPE.CHEVRON
    return MSO_SHAPE.RECTANGLE


def set_fill(shape, element: dict[str, Any]) -> None:
    fill_color = element.get("fill")
    if not fill_color or str(fill_color).lower() == "none":
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if "fill_transparency" in element:
        shape.fill.transparency = float(element["fill_transparency"])


def set_line(line, element: dict[str, Any]) -> None:
    line_color = element.get("stroke") or element.get("color")
    if not line_color or str(line_color).lower() == "none":
        line.fill.background()
        return
    line.color.rgb = rgb(line_color)
    line.width = pt(float(element.get("stroke_width", element.get("width", 1.0))))
    if element.get("dash"):
        from pptx.enum.dml import MSO_LINE_DASH_STYLE

        line.dash_style = MSO_LINE_DASH_STYLE.DASH


def arrow_requested(element: dict[str, Any], keys: tuple[str, ...]) -> bool:
    return any(as_bool(element.get(key), False) for key in keys)


def arrow_type(element: dict[str, Any], keys: tuple[str, ...]) -> str:
    for key in keys:
        value = element.get(key)
        if isinstance(value, str) and value.strip().lower() not in {"", "1", "true", "yes", "y", "on"}:
            return value.strip()
    return str(element.get("arrow_type", "triangle"))


def remove_arrow_child(ln, tag_name: str) -> None:
    for child in list(ln):
        if xml_local_name(child.tag) == tag_name:
            ln.remove(child)


def append_arrow_child(ln, tag_name: str, element: dict[str, Any], keys: tuple[str, ...], prefix: str) -> int:
    from pptx.oxml.xmlchemy import OxmlElement

    remove_arrow_child(ln, tag_name)
    if not arrow_requested(element, keys):
        return 0

    arrow = OxmlElement(f"a:{tag_name}")
    arrow.set("type", arrow_type(element, keys))
    width = element.get(f"{prefix}_arrow_width") or element.get("arrow_width")
    length = element.get(f"{prefix}_arrow_length") or element.get("arrow_length")
    if width:
        arrow.set("w", str(width))
    if length:
        arrow.set("len", str(length))
    ln.append(arrow)
    return 1


def set_arrowheads(line, element: dict[str, Any]) -> int:
    ln = line._get_or_add_ln()
    head_count = append_arrow_child(
        ln,
        "headEnd",
        element,
        ("start_arrow", "begin_arrow", "head_arrow"),
        "start",
    )
    tail_count = append_arrow_child(
        ln,
        "tailEnd",
        element,
        ("end_arrow", "tail_arrow"),
        "end",
    )
    return head_count + tail_count


def add_shape(slide, element: dict[str, Any], counts: dict[str, int]) -> None:
    shape = slide.shapes.add_shape(
        shape_type(element.get("shape")),
        inches(element["x"]),
        inches(element["y"]),
        inches(element["w"]),
        inches(element["h"]),
    )
    shape.name = str(element.get("id", "native_shape"))
    if "radius" in element and hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        shape.adjustments[0] = float(element["radius"])
    clear_inherited_effects(shape, counts)
    set_fill(shape, element)
    set_line(shape.line, element)
    counts["nativeShapes"] += 1


def add_line(slide, element: dict[str, Any], counts: dict[str, int]) -> None:
    from pptx.enum.shapes import MSO_CONNECTOR

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        inches(element["x1"]),
        inches(element["y1"]),
        inches(element["x2"]),
        inches(element["y2"]),
    )
    connector.name = str(element.get("id", "native_line"))
    clear_inherited_effects(connector, counts)
    set_line(connector.line, element)
    arrowhead_count = set_arrowheads(connector.line, element)
    counts["nativeLines"] += 1
    if arrowhead_count:
        counts["nativeArrowLines"] += 1
        counts["nativeArrowheads"] += arrowhead_count
    if element.get("dash"):
        counts["nativeDashedLines"] += 1


def add_text(slide, element: dict[str, Any], counts: dict[str, int]) -> None:
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR

    box = slide.shapes.add_textbox(
        inches(element["x"]),
        inches(element["y"]),
        inches(element["w"]),
        inches(element["h"]),
    )
    box.name = str(element.get("id", "editable_text"))
    if "rotation" in element:
        box.rotation = float(element["rotation"])
    clear_inherited_effects(box, counts)
    frame = box.text_frame
    frame.clear()
    level = trace_level(element)
    frame.word_wrap = as_bool(element.get("word_wrap"), level != "line")
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = inches(float(element.get("margin", 0.03)))
    frame.margin_right = inches(float(element.get("margin", 0.03)))
    frame.margin_top = inches(float(element.get("margin", 0.02)))
    frame.margin_bottom = inches(float(element.get("margin", 0.02)))
    if str(element.get("valign", "top")).lower() == "middle":
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    align = align_map.get(str(element.get("align", "left")).lower(), PP_ALIGN.LEFT)
    lines = str(element.get("text", "")).splitlines() or [""]
    for idx, line_text in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = line_text
        paragraph.alignment = align
        paragraph.level = int(element.get("level", 0))
        if not paragraph.runs:
            run = paragraph.add_run()
            run.text = line_text
        for run in paragraph.runs:
            font = run.font
            font.name = str(element.get("font_face", "Microsoft YaHei"))
            font.size = pt(float(element.get("font_size", 14)))
            font.bold = as_bool(element.get("bold"), False)
            font.italic = as_bool(element.get("italic"), False)
            font.color.rgb = rgb(element.get("color", "#111111"))

    counts["editableTextBodies"] += 1
    if source_kind(element) == "ocr":
        counts["ocrDerivedTextBodies"] += 1
    else:
        counts["manualTextBodies"] += 1
    if level == "line":
        counts["lineTraceTextBodies"] += 1
    else:
        counts["paragraphTraceTextBodies"] += 1
    if as_bool(element.get("semantic_block"), True):
        counts["semanticTextBlocks"] += 1
    if len(lines) > 1:
        counts["mergedMultilineBlocks"] += 1
        counts["splitTextBoxesPrevented"] += 1
    if text_overflow_risk(element, lines):
        counts["textOverflowRiskBodies"] += 1


def add_picture(slide, element: dict[str, Any], counts: dict[str, int], plan_dir: Path) -> None:
    path_value = element.get("path") or element.get("src") or element.get("image")
    if not path_value:
        raise ValueError(f"Picture element missing path/src/image: {element.get('id')}")

    image_path = resolve_resource(str(path_value), plan_dir)
    kwargs = {}
    if "w" in element:
        kwargs["width"] = inches(element["w"])
    if "h" in element:
        kwargs["height"] = inches(element["h"])
    picture = slide.shapes.add_picture(
        str(image_path),
        inches(element.get("x", 0)),
        inches(element.get("y", 0)),
        **kwargs,
    )
    picture.name = str(element.get("id", "cropped_picture"))
    clear_inherited_effects(picture, counts)
    crop = element.get("crop") or {}
    for side in ("left", "right", "top", "bottom"):
        if side in crop:
            setattr(picture, f"crop_{side}", float(crop[side]))
    counts["pictures"] += 1
    if crop:
        counts["croppedPictures"] += 1
    slide_w = float(element.get("slide_width_in", 13.333333))
    slide_h = float(element.get("slide_height_in", 7.5))
    if float(element.get("x", 0)) == 0 and float(element.get("y", 0)) == 0:
        if abs(float(element.get("w", 0)) - slide_w) < 0.02 and abs(float(element.get("h", 0)) - slide_h) < 0.02:
            counts["fullSlidePictures"] += 1


def add_table(slide, element: dict[str, Any], counts: dict[str, int]) -> None:
    rows = element.get("rows") or []
    if not rows:
        raise ValueError(f"Table element has no rows: {element.get('id')}")
    row_count = len(rows)
    col_count = max(len(row) for row in rows)
    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        inches(element["x"]),
        inches(element["y"]),
        inches(element["w"]),
        inches(element["h"]),
    )
    table_shape.name = str(element.get("id", "native_table"))
    clear_inherited_effects(table_shape, counts)
    table = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    font.name = str(element.get("font_face", "Microsoft YaHei"))
                    font.size = pt(float(element.get("font_size", 10)))
                    font.color.rgb = rgb(element.get("color", "#111111"))
    counts["nativeTables"] += 1


def load_plan(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def build(plan: dict[str, Any], out: Path, plan_dir: Path) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation()
    prs.slide_width = inches(plan.get("slide_width_in", 13.333333))
    prs.slide_height = inches(plan.get("slide_height_in", 7.5))

    blank_layout = prs.slide_layouts[6]
    counts = {
        "slides": 0,
        "pictures": 0,
        "fullSlidePictures": 0,
        "nativeShapes": 0,
        "nativeLines": 0,
        "nativeArrowLines": 0,
        "nativeArrowheads": 0,
        "nativeDashedLines": 0,
        "nativeTables": 0,
        "editableTextBodies": 0,
        "semanticTextBlocks": 0,
        "mergedMultilineBlocks": 0,
        "splitTextBoxesPrevented": 0,
        "unwantedSplitTextBoxes": 0,
        "manualTextBodies": 0,
        "ocrDerivedTextBodies": 0,
        "lineTraceTextBodies": 0,
        "paragraphTraceTextBodies": 0,
        "textOverflowRiskBodies": 0,
        "croppedPictures": 0,
        "effectClearedObjects": 0,
        "effectXmlNodesCleared": 0,
    }

    for slide_plan in plan.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        counts["slides"] += 1
        for element in slide_plan.get("elements", []):
            kind = element.get("type")
            if kind == "shape":
                add_shape(slide, element, counts)
            elif kind == "line":
                add_line(slide, element, counts)
            elif kind == "text":
                add_text(slide, element, counts)
            elif kind in {"picture", "image"}:
                add_picture(slide, element, counts, plan_dir)
            elif kind == "table":
                add_table(slide, element, counts)
            else:
                raise ValueError(f"Unsupported element type: {kind}")

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return {
        "status": "generated",
        "pptx": str(out),
        "summary": counts,
        "limitations": plan.get("limitations", []),
    }


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("plan", type=Path)
    parser.add_argument("--out", type=Path, required=True)
    parser.add_argument("--report", type=Path, required=True)
    parser.add_argument("--pptx-deps", type=Path, default=Path(".deps/ppt-to-editable-pptx"))
    args = parser.parse_args()

    add_pptx_deps(args.pptx_deps)
    plan = load_plan(args.plan)
    report = build(plan, args.out, args.plan.resolve().parent)
    args.report.parent.mkdir(parents=True, exist_ok=True)
    args.report.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(report["summary"], ensure_ascii=False))


if __name__ == "__main__":
    main()
