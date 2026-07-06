import importlib.util
import json
import shutil
import tempfile
import unittest
from pathlib import Path


SCRIPT_PATH = (
    Path(__file__).resolve().parents[1]
    / "scripts"
    / "deck_controller.py"
)
SETUP_SCRIPT_PATH = (
    Path(__file__).resolve().parents[1]
    / "scripts"
    / "setup_ocr_runtime.py"
)
DISPATCH_SCRIPT_PATH = (
    Path(__file__).resolve().parents[1]
    / "scripts"
    / "prepare_worker_dispatch.py"
)

spec = importlib.util.spec_from_file_location("deck_controller", SCRIPT_PATH)
deck_controller = importlib.util.module_from_spec(spec)
assert spec.loader is not None
spec.loader.exec_module(deck_controller)

setup_spec = importlib.util.spec_from_file_location("setup_ocr_runtime", SETUP_SCRIPT_PATH)
setup_ocr_runtime = importlib.util.module_from_spec(setup_spec)
assert setup_spec.loader is not None
setup_spec.loader.exec_module(setup_ocr_runtime)


def load_dispatch_module():
    dispatch_spec = importlib.util.spec_from_file_location(
        "prepare_worker_dispatch",
        DISPATCH_SCRIPT_PATH,
    )
    dispatch_module = importlib.util.module_from_spec(dispatch_spec)
    assert dispatch_spec.loader is not None
    dispatch_spec.loader.exec_module(dispatch_module)
    return dispatch_module


TMP_ROOT = Path(tempfile.gettempdir()) / "ppt_to_editable_3_0_integrated_tests"
TMP_ROOT.mkdir(parents=True, exist_ok=True)


def fresh_dir(name: str) -> Path:
    path = TMP_ROOT / name
    if path.exists():
        shutil.rmtree(path)
    path.mkdir(parents=True)
    return path


def write_json(path: Path, data: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def write_json_bom(path: Path, data: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8-sig")


def visual_fidelity_review_fields() -> dict:
    return {
        "font_size_consistency_review": {
            "status": "reviewed",
            "issues": [],
        },
        "complex_visual_strategy_review": {
            "status": "reviewed",
            "unresolved_native_redraw_drift": False,
            "issues": [],
        },
        "text_background_policy": {
            "status": "reviewed",
            "invented_text_backgrounds_added": False,
            "issues": [],
        },
        "non_text_anchor_preservation_review": {
            "status": "reviewed",
            "missing_required_anchors": False,
            "issues": [],
        },
    }


def valid_region_crop_plan() -> list[dict]:
    return [
        {
            "id": "title",
            "strategy": "native_text",
            "source_bbox_px": [60, 40, 1500, 130],
            "reason": "Title text is rebuilt as editable text.",
        },
        {
            "id": "main",
            "strategy": "source_crop",
            "source_bbox_px": [70, 150, 1510, 760],
            "reason": "Main visual region is preserved as source-derived crop when needed.",
        },
        {
            "id": "footer",
            "strategy": "native_text",
            "source_bbox_px": [60, 780, 1500, 850],
            "reason": "Footer text is rebuilt as editable text.",
        },
    ]


def make_image_only_pptx(path: Path, slide_count: int = 2) -> None:
    from PIL import Image
    from pptx import Presentation

    image_path = path.with_suffix(".png")
    Image.new("RGB", (1600, 900), "white").save(image_path)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(slide_count):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    prs.save(path)


def valid_gates_for(
    source_pptx: Path,
    user_choice: str = "all",
    slides: str | list[int] = "",
    conversion_mode: str = "full-deck-convenience",
    user_confirmed_setup: bool = True,
    user_accepted_limited_run: bool = False,
    setup_required: bool = True,
    auto_passed_existing_runtime: bool = False,
    worker_mode: str = "app-native",
    worker_execution_approved: bool = True,
    external_codex_worker_approved: bool = False,
    token_cost_acknowledged: bool = True,
    content_sharing_acknowledged: bool = True,
) -> dict:
    return {
        "schema_version": deck_controller.GATES_SCHEMA_VERSION,
        "source_pptx_sha256": deck_controller.sha256_file(source_pptx),
        "conversion_mode_gate": {
            "user_choice": conversion_mode,
            "user_response_quote": "选择全量省力模式。" if conversion_mode == "full-deck-convenience" else "选择单页省 token 模式。",
        },
        "ocr_runtime_gate": {
            "probe_ocr_runtime_status": deck_controller.OCR_TEXT_USABLE_STATUS
            if auto_passed_existing_runtime
            else deck_controller.OCR_FAILED_STATUS,
            "setup_required": setup_required,
            "auto_passed_existing_runtime": auto_passed_existing_runtime,
            "explained_to_user": setup_required,
            "user_confirmed_setup": user_confirmed_setup,
            "user_accepted_limited_run": user_accepted_limited_run,
            "user_response_quote": "确认，可以继续。" if setup_required else "",
        },
        "scope_and_worker_gate": {
            "user_choice": user_choice,
            "slides": slides,
            "worker_mode": worker_mode,
            "worker_execution_approved": worker_execution_approved,
            "external_codex_worker_approved": external_codex_worker_approved,
            "token_cost_acknowledged": token_cost_acknowledged,
            "content_sharing_acknowledged": content_sharing_acknowledged,
            "user_response_quote": "按这个范围转换。",
        },
    }


class OcrRuntimeContractTests(unittest.TestCase):
    def test_parse_slide_selection_supports_lists_ranges_and_deduping(self) -> None:
        self.assertEqual(deck_controller.parse_slide_selection("1, 3-5, 3"), [1, 3, 4, 5])
        self.assertIsNone(deck_controller.parse_slide_selection(None))
        with self.assertRaises(ValueError):
            deck_controller.parse_slide_selection("0")
        with self.assertRaises(ValueError):
            deck_controller.parse_slide_selection("5-3")
        with self.assertRaises(ValueError):
            deck_controller.parse_slide_selection("a")

    def test_filter_slides_by_selection_uses_requested_scope(self) -> None:
        slides = [{"slide_index": 1}, {"slide_index": 2}, {"slide_index": 3}]
        selected = deck_controller.filter_slides_by_selection(slides, [3, 1])
        self.assertEqual([slide["slide_index"] for slide in selected], [3, 1])
        with self.assertRaises(ValueError):
            deck_controller.filter_slides_by_selection(slides, [4])

    def test_create_run_requires_gates_file_before_initialization(self) -> None:
        run_dir = fresh_dir("scope-gate-required-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")

        with self.assertRaisesRegex(ValueError, "Gates file required"):
            deck_controller.create_run(
                source_pptx,
                run_dir / "output",
                finalize_fallback=False,
                finalize=False,
                all_slides=True,
            )
        self.assertFalse((run_dir / "output").exists())

    def test_probe_deck_returns_read_only_gate_template(self) -> None:
        run_dir = fresh_dir("probe-deck-run")
        source_pptx = run_dir / "source.pptx"
        make_image_only_pptx(source_pptx, slide_count=2)

        original_get_ocr_dependency_status = deck_controller.get_ocr_dependency_status

        def fake_ocr_status(*_args, **_kwargs):
            return {
                "status": "passed",
                "ocr_runtime_status": "passed-text-usable",
                "checks": {},
                "runtime": {},
            }

        deck_controller.get_ocr_dependency_status = fake_ocr_status
        try:
            result = deck_controller.probe_deck(source_pptx)
        finally:
            deck_controller.get_ocr_dependency_status = original_get_ocr_dependency_status

        self.assertEqual(result["schema_version"], deck_controller.PROBE_SCHEMA_VERSION)
        self.assertEqual(result["slide_count"], 2)
        self.assertTrue(result["all_slides_are_full_page_images"])
        self.assertEqual(result["gates_file_template"]["schema_version"], deck_controller.GATES_SCHEMA_VERSION)
        self.assertIn("conversion_mode_gate", result["gates_file_template"])
        self.assertIn("conversion_mode_gate", result["user_confirmation_required"])
        self.assertIn("scope_and_worker_gate", result["gates_file_template"])
        self.assertIn("scope_and_worker_gate", result["user_confirmation_required"])
        self.assertNotIn("scope_gate", result["user_confirmation_required"])
        self.assertFalse(
            result["user_confirmation_required"]["ocr_runtime_gate"]["requires_user_confirmation"]
        )
        mode_question = result["user_confirmation_required"]["conversion_mode_gate"]["question"]
        self.assertIn("单页省 token 模式", mode_question)
        self.assertIn("多页多 Agent 高质量模式", mode_question)
        self.assertNotIn("多页顺序省 token 模式", mode_question)
        self.assertIn("发单张 PNG", mode_question)
        self.assertIn("省 token", mode_question)
        conversion_options = result["user_confirmation_required"]["conversion_mode_gate"]["options"]
        self.assertEqual(
            conversion_options,
            [
                "single-page-token-saving",
                "multi-agent-high-quality",
            ],
        )
        for forbidden in (
            "app-native",
            "external-codex",
            "worker runtime",
            "Codex worker",
            "worker_mode",
            "scope_and_worker_gate",
        ):
            self.assertNotIn(forbidden, mode_question)
        scope_question = result["user_confirmation_required"]["scope_and_worker_gate"]["question"]
        self.assertIn("这份 PPT 一共有 2 页", scope_question)
        self.assertIn("先转换 1 页看看效果", scope_question)
        self.assertIn("消耗的 token 越多", scope_question)
        self.assertNotIn("同一个转换流程里逐页处理", scope_question)
        self.assertIn("Pro 订阅用户", scope_question)
        self.assertIn("仅限 Codex Pro", scope_question)
        self.assertIn("非常消耗 token", scope_question)
        self.assertIn("独立转换任务", scope_question)
        self.assertIn("读取对应页面的图片和文字内容", scope_question)
        for forbidden in (
            "app-native",
            "external-codex",
            "worker runtime",
            "Codex worker",
            "worker_mode",
            "scope_and_worker_gate",
        ):
            self.assertNotIn(forbidden, scope_question)
        self.assertFalse((run_dir / "source_slides").exists())

    def test_conversion_mode_gate_rejects_retired_sequential_mode(self) -> None:
        run_dir = fresh_dir("gates-rejects-retired-sequential-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        write_json(
            gates_file,
            valid_gates_for(
                source_pptx,
                conversion_mode="multi-page-sequential-token-saving",
                worker_mode="sequential",
                content_sharing_acknowledged=False,
                setup_required=False,
                auto_passed_existing_runtime=True,
            ),
        )

        with self.assertRaisesRegex(ValueError, "retired|下线|not supported"):
            deck_controller.load_and_validate_gates(
                gates_file,
                source_pptx,
                slides=None,
                all_slides=True,
                allow_ocr_not_ready=False,
            )

    def test_conversion_mode_gate_rejects_single_page_mode_with_all_slides(self) -> None:
        run_dir = fresh_dir("gates-mode-rejects-all-slides-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        write_json(
            gates_file,
            valid_gates_for(
                source_pptx,
                conversion_mode="single-page-token-saving",
                user_choice="all",
                slides="",
                setup_required=False,
                auto_passed_existing_runtime=True,
            ),
        )

        with self.assertRaisesRegex(ValueError, "Conversion Mode Gate"):
            deck_controller.load_and_validate_gates(
                gates_file,
                source_pptx,
                slides=None,
                all_slides=True,
                allow_ocr_not_ready=False,
            )

    def test_conversion_mode_gate_allows_single_page_mode_with_one_sample_slide(self) -> None:
        run_dir = fresh_dir("gates-mode-allows-sample-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        write_json(
            gates_file,
            valid_gates_for(
                source_pptx,
                conversion_mode="single-page-token-saving",
                user_choice="sample",
                slides="2",
                setup_required=False,
                auto_passed_existing_runtime=True,
            ),
        )

        gates, _ = deck_controller.load_and_validate_gates(
            gates_file,
            source_pptx,
            slides="2",
            all_slides=False,
            allow_ocr_not_ready=False,
        )
        self.assertEqual(
            gates["conversion_mode_gate"]["user_choice"],
            "single-page-token-saving",
        )

    def test_load_and_validate_gates_allows_ready_ocr_without_setup_confirmation(self) -> None:
        run_dir = fresh_dir("gates-ready-ocr-no-setup-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        write_json(
            gates_file,
            valid_gates_for(
                source_pptx,
                setup_required=False,
                auto_passed_existing_runtime=True,
                user_confirmed_setup=False,
            ),
        )

        gates, _ = deck_controller.load_and_validate_gates(
            gates_file,
            source_pptx,
            slides=None,
            all_slides=True,
            allow_ocr_not_ready=False,
        )
        self.assertTrue(gates["ocr_runtime_gate"]["auto_passed_existing_runtime"])

    def test_load_and_validate_gates_rejects_scope_without_worker_authorization(self) -> None:
        run_dir = fresh_dir("gates-rejects-legacy-scope-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        gates = {
            "schema_version": deck_controller.GATES_SCHEMA_VERSION,
            "source_pptx_sha256": deck_controller.sha256_file(source_pptx),
            "conversion_mode_gate": {
                "user_choice": "full-deck-convenience",
                "user_response_quote": "选择全量省力模式。",
            },
            "ocr_runtime_gate": {
                "setup_required": False,
                "auto_passed_existing_runtime": True,
                "probe_ocr_runtime_status": deck_controller.OCR_TEXT_USABLE_STATUS,
            },
            "scope_gate": {
                "user_choice": "all",
                "slides": "",
                "user_response_quote": "全部转换。",
            },
        }
        write_json(gates_file, gates)

        with self.assertRaisesRegex(ValueError, "scope_and_worker_gate"):
            deck_controller.load_and_validate_gates(
                gates_file,
                source_pptx,
                slides=None,
                all_slides=True,
                allow_ocr_not_ready=False,
            )

    def test_external_codex_worker_mode_requires_explicit_external_approval(self) -> None:
        run_dir = fresh_dir("gates-external-worker-approval-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        write_json(
            gates_file,
            valid_gates_for(
                source_pptx,
                setup_required=False,
                auto_passed_existing_runtime=True,
                worker_mode="external-codex",
                external_codex_worker_approved=False,
            ),
        )

        with self.assertRaisesRegex(ValueError, "external_codex_worker_approved"):
            deck_controller.load_and_validate_gates(
                gates_file,
                source_pptx,
                slides=None,
                all_slides=True,
                allow_ocr_not_ready=False,
            )

    def test_load_and_validate_gates_enforces_scope_match(self) -> None:
        run_dir = fresh_dir("gates-scope-match-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        write_json(gates_file, valid_gates_for(source_pptx, user_choice="selected", slides="1,3"))

        with self.assertRaisesRegex(ValueError, "Scope \\+ Worker Gate mismatch"):
            deck_controller.load_and_validate_gates(
                gates_file,
                source_pptx,
                slides="1,2",
                all_slides=False,
                allow_ocr_not_ready=False,
            )

    def test_create_run_records_user_confirmation_gates(self) -> None:
        run_dir = fresh_dir("create-run-records-gates")
        source_pptx = run_dir / "source.pptx"
        make_image_only_pptx(source_pptx, slide_count=1)
        gates_file = run_dir / "gates.json"
        write_json(gates_file, valid_gates_for(source_pptx))
        original_get_ocr_dependency_status = deck_controller.get_ocr_dependency_status

        def fake_ocr_status(*_args, **_kwargs):
            return {
                "status": "passed",
                "ocr_runtime_status": "passed-text-usable",
                "checks": {},
                "runtime": {},
            }

        deck_controller.get_ocr_dependency_status = fake_ocr_status
        try:
            result = deck_controller.create_run(
                source_pptx,
                run_dir / "output",
                finalize_fallback=False,
                finalize=False,
                all_slides=True,
                gates_file=gates_file,
            )
        finally:
            deck_controller.get_ocr_dependency_status = original_get_ocr_dependency_status

        manifest = json.loads(Path(result["deck_manifest"]).read_text(encoding="utf-8"))
        self.assertEqual(result["user_confirmation_gates_status"], "passed")
        self.assertEqual(manifest["user_confirmation_gates"]["status"], "passed")
        self.assertEqual(
            manifest["user_confirmation_gates"]["conversion_mode_gate"]["user_choice"],
            "full-deck-convenience",
        )
        self.assertEqual(
            manifest["user_confirmation_gates"]["scope_and_worker_gate"]["user_choice"],
            "all",
        )
        self.assertTrue(
            manifest["user_confirmation_gates"]["scope_and_worker_gate"]["worker_execution_approved"]
        )

    def test_prepare_worker_dispatch_blocks_unapproved_external_worker(self) -> None:
        run_dir = fresh_dir("prepare-dispatch-blocks-external-run")
        source_pptx = run_dir / "source.pptx"
        make_image_only_pptx(source_pptx, slide_count=1)
        gates_file = run_dir / "gates.json"
        write_json(
            gates_file,
            valid_gates_for(
                source_pptx,
                setup_required=False,
                auto_passed_existing_runtime=True,
                worker_mode="manual",
                external_codex_worker_approved=False,
            ),
        )
        original_get_ocr_dependency_status = deck_controller.get_ocr_dependency_status

        def fake_ocr_status(*_args, **_kwargs):
            return {
                "status": "passed",
                "ocr_runtime_status": "passed-text-usable",
                "checks": {},
                "runtime": {},
            }

        deck_controller.get_ocr_dependency_status = fake_ocr_status
        try:
            deck_controller.create_run(
                source_pptx,
                run_dir / "output",
                finalize_fallback=False,
                finalize=False,
                all_slides=True,
                gates_file=gates_file,
            )
        finally:
            deck_controller.get_ocr_dependency_status = original_get_ocr_dependency_status

        dispatcher = load_dispatch_module()
        report = dispatcher.build_dispatch_report(
            run_dir / "output",
            worker_mode="external-codex",
        )
        self.assertEqual(report["status"], "blocked-external-worker-not-approved")
        self.assertEqual(report["conversion_mode"], "full-deck-convenience")
        self.assertFalse(report["automatic_external_codex_allowed"])

    def test_setup_ocr_runtime_yes_requires_confirmed_gates(self) -> None:
        run_dir = fresh_dir("setup-ocr-gates-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        write_json(
            gates_file,
            valid_gates_for(
                source_pptx,
                user_confirmed_setup=False,
            ),
        )

        with self.assertRaisesRegex(ValueError, "user_confirmed_setup=true"):
            setup_ocr_runtime.load_and_validate_setup_gates(gates_file)

    def test_ocr_runtime_gate_rejects_quality_manifest_when_text_is_not_usable(self) -> None:
        run_dir = fresh_dir("ocr-gate-required-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        original_get_ocr_dependency_status = deck_controller.get_ocr_dependency_status

        def fake_ocr_status(*_args, **_kwargs):
            return {
                "status": "warning",
                "ocr_runtime_status": "failed",
                "checks": {},
                "runtime": {},
            }

        deck_controller.get_ocr_dependency_status = fake_ocr_status
        try:
            with self.assertRaisesRegex(RuntimeError, "OCR Runtime Gate required"):
                deck_controller.build_deck_manifest(
                    source_pptx,
                    run_dir,
                    "run-1",
                    {"width_emu": 1, "height_emu": 1},
                    [
                        {
                            "slide_index": 1,
                            "is_single_full_slide_picture": True,
                            "pictures": [{"image_format": "PNG"}],
                        }
                    ],
                )
        finally:
            deck_controller.get_ocr_dependency_status = original_get_ocr_dependency_status

    def test_finalize_reuses_existing_manifest_without_rebuilding_preflight(self) -> None:
        run_dir = fresh_dir("finalize-reuses-existing-manifest-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        write_json(
            run_dir / "deck_manifest.json",
            {
                "slide_count": 1,
                "preflight": {
                    "status": "passed",
                    "ocr_dependency_status": "passed",
                    "ocr_runtime_status": "passed-text-usable",
                    "all_slides_are_full_page_images": True,
                    "ocr_runtime": {
                        "mode": "external",
                        "python": "external-ocr-python",
                    },
                },
                "output": {
                    "final_deck": str(run_dir / "output" / "final-deck.pptx"),
                    "deck_qa_report": str(run_dir / "output" / "deck-level-qa-report.json"),
                },
                "slides": [],
            },
        )

        original_extract = deck_controller.extract_image_only_slides
        original_assemble = deck_controller.assemble_final_deck_from_results
        captured = {}

        def fail_extract(*_args, **_kwargs):
            raise AssertionError("finalize must not rebuild slide manifests or OCR preflight")

        def fake_assemble(manifest, _run_dir):
            captured["manifest"] = manifest
            manifest["output"]["final_deck"] = str(run_dir / "output" / "final-deck.pptx")
            manifest["output"]["deck_qa_report"] = str(run_dir / "output" / "deck-level-qa-report.json")
            return {"final_deck": manifest["output"]["final_deck"]}

        deck_controller.extract_image_only_slides = fail_extract
        deck_controller.assemble_final_deck_from_results = fake_assemble
        try:
            result = deck_controller.create_run(
                source_pptx,
                run_dir,
                finalize_fallback=False,
                finalize=True,
            )
        finally:
            deck_controller.extract_image_only_slides = original_extract
            deck_controller.assemble_final_deck_from_results = original_assemble

        self.assertEqual(result["ocr_dependency_status"], "passed")
        self.assertEqual(result["ocr_runtime_status"], "passed-text-usable")
        self.assertEqual(captured["manifest"]["preflight"]["ocr_runtime"]["python"], "external-ocr-python")

    def test_success_is_rejected_when_ocr_text_is_geometry_only(self) -> None:
        run_dir = fresh_dir("geometry-only-run")
        attempt_dir = run_dir / "slide_jobs" / "slide_01" / "attempt_01"
        attempt_dir.mkdir(parents=True)
        (attempt_dir / "slide.pptx").write_bytes(b"placeholder")
        write_json(
            attempt_dir / "progress.json",
            {
                "history": [
                    {"status": "started"},
                    {"status": "ocr_done"},
                    {"status": "plan_done"},
                    {"status": "packaging_done"},
                    {"status": "qa_done"},
                ]
            },
        )
        write_json(
            attempt_dir / "qa_summary.json",
            {
                "status": "passed",
                "mode": "mixed-reconstruction",
                "references_read": [
                    {"id": "single_page_engine_skill_3_0"},
                    {"id": "slide_job_manifest"},
                ],
                "font_policy_applied": "microsoft-yahei-all-runs",
                **visual_fidelity_review_fields(),
                "effects_cleared": True,
                "rendered_png_reviewed": True,
                "region_crop_plan": valid_region_crop_plan(),
                "editable_text_count": 8,
                "native_shapes_count": 4,
                "source_crop_count": 3,
                "visual_qa": {"known_visual_risks": []},
                "failure_reason": None,
                "limitations": [],
                "recommended_next_round_change": "",
                "worker_identity": {
                    "worker_thread_id": "worker-1",
                    "execution_model": "delegated-real-per-slide-worker",
                    "slide_scope": [1],
                },
                "ocr_runtime_status": "passed-geometry-only",
                "single_page_engine_used": "ppt-to-editable-v3-1-preview",
                "single_page_engine_scripts_used": ["run_reconstruction_qa.py"],
            },
        )

        result = deck_controller.collect_slide_result(
            {
                "slide_index": 1,
                "required_reading": [
                    {"id": "single_page_engine_skill_3_0", "must_read": True},
                    {"id": "slide_job_manifest", "must_read": True},
                ],
                "qa_summary_required_fields": [
                    "status",
                    "mode",
                    "references_read",
                    "font_policy_applied",
                    "font_size_consistency_review",
                    "effects_cleared",
                    "rendered_png_reviewed",
                    "region_crop_plan",
                    "complex_visual_strategy_review",
                    "text_background_policy",
                    "editable_text_count",
                    "native_shapes_count",
                    "source_crop_count",
                    "visual_qa",
                    "failure_reason",
                    "limitations",
                    "recommended_next_round_change",
                    "worker_identity",
                    "ocr_runtime_status",
                ],
                "required_progress_statuses": deck_controller.REQUIRED_PROGRESS_STATUSES,
                "ocr_dependency_status_at_dispatch": "passed",
                "ocr_runtime_status_at_dispatch": "passed-geometry-only",
            },
            run_dir,
        )

        self.assertEqual(result["status"], "failed-retryable")
        self.assertIn("ocr_text_not_usable", result["failure_reason"])

    def test_manifest_requires_worker_to_report_ocr_runtime_status(self) -> None:
        run_dir = fresh_dir("manifest-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")
        manifest = deck_controller.make_slide_job_manifest(
            "run-1",
            {
                "slide_index": 1,
                "pictures": [
                    {
                        "extracted_path": str(source),
                        "image_size": {"width_px": 1600, "height_px": 900},
                    }
                ],
            },
            run_dir,
            ocr_runtime_status="passed-text-usable",
        )

        self.assertIn("ocr_runtime_status", manifest["qa_summary_required_fields"])
        self.assertEqual(manifest["ocr_runtime"]["status"], "passed-text-usable")
        self.assertEqual(manifest["worker_contract"]["reuse_skill"], "ppt-to-editable-v3-1-preview")
        required_ids = {item["id"] for item in manifest["required_reading"] if item["must_read"]}
        self.assertIn("worker_compact_protocol", required_ids)
        self.assertIn("page_conversion_contract", required_ids)
        self.assertIn("crop_manifest_contract", required_ids)
        fallback_ids = {
            item["id"]
            for item in manifest["fallback_reading"]
            if item.get("read_when_blocked_or_needed")
        }
        self.assertIn("single_page_engine_skill_3_1", fallback_ids)
        self.assertIn("single_page_engine_reconstruction_qa", fallback_ids)
        self.assertIn("single_page_engine_source_crop_tools", fallback_ids)
        self.assertNotIn("active_skill", required_ids)
        self.assertIn("single_page_engine_used", manifest["qa_summary_required_fields"])
        self.assertIn("single_page_engine_scripts_used", manifest["qa_summary_required_fields"])
        self.assertIn("fallback_references_read", manifest["qa_summary_required_fields"])
        task_text = (run_dir / "slide_jobs" / "slide_01" / "AGENT_TASK.md").read_text(encoding="utf-8")
        self.assertIn("ppt-to-editable", task_text)
        self.assertNotIn("ppt-to-editable-2-0", task_text)
        self.assertIn("run_reconstruction_qa.py", task_text)
        self.assertIn("write_slide_qa_summary_from_single_page.py", task_text)
        self.assertIn("Do not force native reconstruction for complex shapes", task_text)
        self.assertIn("textless crop plus editable text", task_text)
        self.assertTrue(manifest["quality_rules"]["complex_shapes_should_use_textless_crop_when_native_redraw_would_drift"])

    def test_manifest_and_agent_task_include_round_23_visual_fidelity_guards(self) -> None:
        run_dir = fresh_dir("visual-fidelity-guards-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")
        manifest = deck_controller.make_slide_job_manifest(
            "run-1",
            {
                "slide_index": 1,
                "pictures": [
                    {
                        "extracted_path": str(source),
                        "image_size": {"width_px": 1600, "height_px": 900},
                    }
                ],
            },
            run_dir,
            ocr_runtime_status="passed-text-usable",
        )

        quality_rules = manifest["quality_rules"]
        self.assertTrue(quality_rules["must_preserve_relative_font_size_consistency"])
        self.assertTrue(quality_rules["must_not_shrink_text_to_fit_when_source_scale_would_be_lost"])
        self.assertTrue(quality_rules["must_use_crop_first_for_complex_artwork_regions"])
        self.assertTrue(quality_rules["must_not_add_invented_text_background_fills"])
        self.assertIn("font_size_consistency_review", manifest["qa_summary_required_fields"])
        self.assertIn("complex_visual_strategy_review", manifest["qa_summary_required_fields"])
        self.assertIn("text_background_policy", manifest["qa_summary_required_fields"])

        task_text = (run_dir / "slide_jobs" / "slide_01" / "AGENT_TASK.md").read_text(encoding="utf-8")
        self.assertIn("Do not shrink body text into tiny outliers", task_text)
        self.assertIn("crop-first", task_text)
        self.assertIn("Do not add pale rectangles or invented backing fills behind editable text", task_text)

    def test_manifest_and_agent_task_include_round_24_non_text_anchor_guards(self) -> None:
        run_dir = fresh_dir("non-text-anchor-guards-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")
        manifest = deck_controller.make_slide_job_manifest(
            "run-1",
            {
                "slide_index": 1,
                "pictures": [
                    {
                        "extracted_path": str(source),
                        "image_size": {"width_px": 1600, "height_px": 900},
                    }
                ],
            },
            run_dir,
            ocr_runtime_status="passed-text-usable",
        )

        quality_rules = manifest["quality_rules"]
        self.assertTrue(quality_rules["must_preserve_non_text_visual_anchors"])
        self.assertTrue(quality_rules["must_not_erase_arrows_connectors_step_markers_or_card_edges"])
        self.assertIn("non_text_anchor_preservation_review", manifest["qa_summary_required_fields"])

        task_text = (run_dir / "slide_jobs" / "slide_01" / "AGENT_TASK.md").read_text(encoding="utf-8")
        self.assertIn("Do not erase non-text visual anchors", task_text)
        self.assertIn("arrows, arrowheads, route lines, connectors, step markers", task_text)
        self.assertIn("Record `non_text_anchor_preservation_review`", task_text)

    def test_manifest_and_agent_task_include_p1_execution_contracts(self) -> None:
        run_dir = fresh_dir("p1-execution-contracts-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")
        manifest = deck_controller.make_slide_job_manifest(
            "run-1",
            {
                "slide_index": 1,
                "pictures": [
                    {
                        "extracted_path": str(source),
                        "image_size": {"width_px": 1600, "height_px": 900},
                    }
                ],
            },
            run_dir,
            ocr_runtime_status="passed-text-usable",
        )

        slide_dir = run_dir / "slide_jobs" / "slide_01"
        self.assertTrue((slide_dir / "page_conversion_contract.json").exists())
        self.assertTrue((slide_dir / "crop_manifest_contract.json").exists())

        page_contract = manifest["page_conversion_contract"]
        self.assertEqual(page_contract["default_font_family"], "Microsoft YaHei")
        self.assertEqual(page_contract["fallback_decision"], "prefer_fidelity_over_extra_native_shapes")
        self.assertTrue(page_contract["text_rules"]["meaningful_text_native_when_practical"])
        self.assertTrue(page_contract["shape_rules"]["clear_unwanted_effects"])
        self.assertTrue(page_contract["visual_rules"]["complex_visuals_crop_first"])

        crop_contract = manifest["crop_manifest_contract"]
        self.assertEqual(crop_contract["artifact"], "crop_manifest.json")
        self.assertIn("textless_crop_plus_editable_text", crop_contract["allowed_region_strategies"])
        self.assertIn("must_not_erase_non_text_anchors", crop_contract["review_gates"])
        self.assertIn("crop_manifest", manifest["qa_summary_required_fields"])

        route_ids = {row["id"] for row in manifest["routing_decision_table"]}
        self.assertIn("pure_text_to_native_text", route_ids)
        self.assertIn("complex_visual_to_source_crop", route_ids)
        self.assertIn("text_in_complex_visual_to_textless_crop_overlay", route_ids)
        self.assertIn("ocr_unusable_to_fallback", route_ids)

        policy = manifest["reference_loading_policy"]
        self.assertIn("slide_job_manifest", policy["must_read_ids"])
        self.assertIn("ocr_troubleshooting", policy["conditional_reading_ids"])
        self.assertIn("old_round_outputs", policy["forbidden_context"])
        self.assertIn("historical_test_reports", policy["forbidden_context"])

        task_text = (slide_dir / "AGENT_TASK.md").read_text(encoding="utf-8")
        self.assertIn("page_conversion_contract.json", task_text)
        self.assertIn("crop_manifest_contract.json", task_text)
        self.assertIn("Routing Decision Table", task_text)
        self.assertIn("Do not read historical test reports", task_text)

    def test_collect_slide_result_accepts_utf8_bom_qa_summary(self) -> None:
        run_dir = fresh_dir("bom-qa-summary-run")
        attempt_dir = run_dir / "slide_jobs" / "slide_01" / "attempt_01"
        render_dir = attempt_dir / "render"
        render_dir.mkdir(parents=True)
        rendered_png = render_dir / "slide_01.png"
        rendered_png.write_bytes(b"png")
        (attempt_dir / "slide.pptx").write_bytes(b"placeholder")
        write_json(
            attempt_dir / "progress.json",
            {
                "history": [
                    {"status": "started"},
                    {"status": "ocr_done"},
                    {"status": "plan_done"},
                    {"status": "packaging_done"},
                    {"status": "qa_done"},
                ]
            },
        )
        write_json(
            attempt_dir / "render_report.json",
            {
                "status": "rendered",
                "rendered_png": str(rendered_png),
            },
        )
        write_json_bom(
            attempt_dir / "qa_summary.json",
            {
                "status": "passed",
                "mode": "mixed-reconstruction",
                "references_read": [
                    {"id": "single_page_engine_skill_3_0"},
                    {"id": "single_page_engine_source_crop_tools"},
                    {"id": "single_page_engine_reconstruction_qa"},
                    {"id": "slide_job_manifest"},
                ],
                "font_policy_applied": "microsoft-yahei-all-runs",
                **visual_fidelity_review_fields(),
                "effects_cleared": True,
                "rendered_png_reviewed": True,
                "render_report": "render_report.json",
                "region_crop_plan": valid_region_crop_plan(),
                "editable_text_count": 8,
                "native_shapes_count": 4,
                "source_crop_count": 3,
                "visual_qa": {"known_visual_risks": []},
                "failure_reason": None,
                "limitations": [],
                "recommended_next_round_change": "",
                "worker_identity": {
                    "worker_thread_id": "worker-bom-1",
                    "execution_model": "delegated-real-per-slide-worker",
                    "slide_scope": [1],
                },
                "ocr_runtime_status": "passed-text-usable",
                "single_page_engine_used": "ppt-to-editable-v3-1-preview",
                "single_page_engine_scripts_used": ["package_reconstruction_deck.py"],
            },
        )

        result = deck_controller.collect_slide_result(
            {
                "slide_index": 1,
                "required_reading": [
                    {"id": "single_page_engine_skill_3_0", "must_read": True},
                    {"id": "single_page_engine_source_crop_tools", "must_read": True},
                    {"id": "single_page_engine_reconstruction_qa", "must_read": True},
                    {"id": "slide_job_manifest", "must_read": True},
                ],
                "qa_summary_required_fields": [
                    "status",
                    "mode",
                    "references_read",
                    "font_policy_applied",
                    "font_size_consistency_review",
                    "effects_cleared",
                    "rendered_png_reviewed",
                    "render_report",
                    "region_crop_plan",
                    "complex_visual_strategy_review",
                    "text_background_policy",
                    "editable_text_count",
                    "native_shapes_count",
                    "source_crop_count",
                    "visual_qa",
                    "failure_reason",
                    "limitations",
                    "recommended_next_round_change",
                    "worker_identity",
                    "ocr_runtime_status",
                    "single_page_engine_used",
                    "single_page_engine_scripts_used",
                ],
                "required_progress_statuses": deck_controller.REQUIRED_PROGRESS_STATUSES,
                "ocr_dependency_status_at_dispatch": "passed",
                "ocr_runtime_status_at_dispatch": "passed-text-usable",
            },
            run_dir,
        )

        self.assertEqual(result["status"], "passed", result.get("failure_reason"))

    def test_collect_slide_result_auto_closes_raw_single_page_summary(self) -> None:
        run_dir = fresh_dir("auto-close-single-page-summary-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")
        manifest = deck_controller.make_slide_job_manifest(
            "run-1",
            {
                "slide_index": 1,
                "pictures": [
                    {
                        "extracted_path": str(source),
                        "image_size": {"width_px": 1600, "height_px": 900},
                    }
                ],
            },
            run_dir,
            ocr_runtime_status="passed-text-usable",
        )
        slide_dir = run_dir / "slide_jobs" / "slide_01"
        attempt_dir = slide_dir / "attempt_01"
        rendered_png = attempt_dir / "render" / "slide_01.png"
        rendered_png.parent.mkdir(parents=True, exist_ok=True)
        rendered_png.write_bytes(b"png")
        (attempt_dir / "slide.pptx").write_bytes(b"placeholder")
        (attempt_dir / "visual-qa-comparison.html").write_text("<html></html>", encoding="utf-8")
        write_json(
            attempt_dir / "progress.json",
            {
                "history": [
                    {"status": "started"},
                    {"status": "ocr_done"},
                    {"status": "plan_done"},
                    {"status": "packaging_done"},
                ]
            },
        )
        write_json(
            attempt_dir / "source_reconstruction_plan.json",
            {
                "mode": "mixed_reconstruction",
                "strategy_review": {
                    "regions": [
                        {
                            "id": "header",
                            "actual_strategy": "native_text",
                            "source_bbox_px": [60, 40, 1500, 130],
                            "reason": "Header text is rebuilt as editable text.",
                        },
                        {
                            "id": "body",
                            "actual_strategy": "textless_crop_plus_editable_text",
                            "source_bbox_px": [70, 150, 1510, 760],
                            "reason": "Main visual uses textless crop plus editable text.",
                        },
                        {
                            "id": "footer",
                            "actual_strategy": "native_text",
                            "source_bbox_px": [60, 780, 1500, 850],
                            "reason": "Footer text is rebuilt as editable text.",
                        },
                    ]
                },
            },
        )
        write_json(
            attempt_dir / "packaging_report.json",
            {
                "summary": {
                    "editableTextBodies": 12,
                    "nativeShapes": 7,
                    "effectClearedObjects": 9,
                }
            },
        )
        write_json(
            attempt_dir / "editability_report.json",
            {"summary": {"editableTextBodies": 12, "nativeShapes": 7}},
        )
        write_json(
            attempt_dir / "crop_report.json",
            {"outputs": [{"id": "region-a"}, {"id": "region-b"}, {"id": "region-c"}]},
        )
        write_json_bom(
            attempt_dir / "render_report.json",
            {
                "status": "rendered",
                "slides": [{"index": 1, "png": str(rendered_png), "exists": True}],
            },
        )
        write_json(
            attempt_dir / "qa_summary.json",
            {
                "status": "pass",
                "steps": [
                    {"name": "package", "command": ["python", "package_reconstruction_deck.py"]},
                    {"name": "inspect_editability", "command": ["python", "inspect_pptx_editability.py"]},
                ],
                "limitations": [],
            },
        )

        result = deck_controller.collect_slide_result(
            {
                "slide_index": 1,
                "required_reading": manifest["required_reading"],
                "qa_summary_required_fields": manifest["qa_summary_required_fields"],
                "required_progress_statuses": deck_controller.REQUIRED_PROGRESS_STATUSES,
                "ocr_dependency_status_at_dispatch": "passed",
                "ocr_runtime_status_at_dispatch": "passed-text-usable",
            },
            run_dir,
        )

        self.assertEqual(result["status"], "passed", result.get("failure_reason"))
        self.assertTrue((attempt_dir / "qa_summary.single-page.json").exists())
        summary = json.loads((attempt_dir / "qa_summary.json").read_text(encoding="utf-8"))
        self.assertEqual(summary["single_page_engine_used"], "ppt-to-editable-v3-1-preview")
        self.assertEqual(summary["worker_identity"]["worker_agent_id"], "controller-autoclosed-slide-01")

    def test_success_is_rejected_when_render_report_is_missing(self) -> None:
        run_dir = fresh_dir("missing-render-report-run")
        attempt_dir = run_dir / "slide_jobs" / "slide_01" / "attempt_01"
        attempt_dir.mkdir(parents=True)
        (attempt_dir / "slide.pptx").write_bytes(b"placeholder")
        write_json(
            attempt_dir / "progress.json",
            {
                "history": [
                    {"status": "started"},
                    {"status": "ocr_done"},
                    {"status": "plan_done"},
                    {"status": "packaging_done"},
                    {"status": "qa_done"},
                ]
            },
        )
        write_json(
            attempt_dir / "qa_summary.json",
            {
                "status": "passed",
                "mode": "mixed-reconstruction",
                "references_read": [
                    {"id": "single_page_engine_skill_3_0"},
                    {"id": "slide_job_manifest"},
                ],
                "font_policy_applied": "microsoft-yahei-all-runs",
                **visual_fidelity_review_fields(),
                "effects_cleared": True,
                "rendered_png_reviewed": True,
                "region_crop_plan": valid_region_crop_plan(),
                "editable_text_count": 8,
                "native_shapes_count": 4,
                "source_crop_count": 3,
                "visual_qa": {"known_visual_risks": []},
                "failure_reason": None,
                "limitations": [],
                "recommended_next_round_change": "",
                "worker_identity": {
                    "worker_thread_id": "worker-1",
                    "execution_model": "delegated-real-per-slide-worker",
                    "slide_scope": [1],
                },
                "ocr_runtime_status": "passed-text-usable",
                "render_report": "render_report.json",
                "single_page_engine_used": "ppt-to-editable-v3-1-preview",
                "single_page_engine_scripts_used": ["run_reconstruction_qa.py"],
            },
        )

        result = deck_controller.collect_slide_result(
            {
                "slide_index": 1,
                "required_reading": [
                    {"id": "single_page_engine_skill_3_0", "must_read": True},
                    {"id": "slide_job_manifest", "must_read": True},
                ],
                "qa_summary_required_fields": [
                    "status",
                    "mode",
                    "references_read",
                    "font_policy_applied",
                    "font_size_consistency_review",
                    "effects_cleared",
                    "rendered_png_reviewed",
                    "region_crop_plan",
                    "complex_visual_strategy_review",
                    "text_background_policy",
                    "editable_text_count",
                    "native_shapes_count",
                    "source_crop_count",
                    "visual_qa",
                    "failure_reason",
                    "limitations",
                    "recommended_next_round_change",
                    "worker_identity",
                    "ocr_runtime_status",
                    "render_report",
                ],
                "required_progress_statuses": deck_controller.REQUIRED_PROGRESS_STATUSES,
                "ocr_dependency_status_at_dispatch": "passed",
                "ocr_runtime_status_at_dispatch": "passed-text-usable",
            },
            run_dir,
        )

        self.assertEqual(result["status"], "failed-retryable")
        self.assertIn("render_report", result["failure_reason"])

    def test_success_is_rejected_when_invented_text_background_is_reported(self) -> None:
        run_dir = fresh_dir("invented-text-background-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")
        manifest = deck_controller.make_slide_job_manifest(
            "run-1",
            {
                "slide_index": 1,
                "pictures": [
                    {
                        "extracted_path": str(source),
                        "image_size": {"width_px": 1600, "height_px": 900},
                    }
                ],
            },
            run_dir,
            ocr_runtime_status="passed-text-usable",
        )
        attempt_dir = run_dir / "slide_jobs" / "slide_01" / "attempt_01"
        rendered_png = attempt_dir / "render" / "slide_01.png"
        rendered_png.parent.mkdir(parents=True, exist_ok=True)
        rendered_png.write_bytes(b"png")
        (attempt_dir / "slide.pptx").write_bytes(b"placeholder")
        write_json(
            attempt_dir / "progress.json",
            {
                "history": [
                    {"status": "started"},
                    {"status": "ocr_done"},
                    {"status": "plan_done"},
                    {"status": "packaging_done"},
                    {"status": "qa_done"},
                ]
            },
        )
        write_json(
            attempt_dir / "render_report.json",
            {
                "status": "rendered",
                "rendered_png": str(rendered_png),
            },
        )
        review_fields = visual_fidelity_review_fields()
        review_fields["text_background_policy"] = {
            "status": "failed",
            "invented_text_backgrounds_added": True,
            "issues": [{"issue": "invented-text-background", "element_id": "body-backing"}],
        }
        write_json(
            attempt_dir / "qa_summary.json",
            {
                "status": "passed",
                "mode": "mixed-reconstruction",
                "references_read": [{"id": item["id"]} for item in manifest["required_reading"]],
                "fallback_references_read": [],
                "font_policy_applied": "microsoft-yahei-all-runs",
                **review_fields,
                "effects_cleared": True,
                "rendered_png_reviewed": True,
                "render_report": str(attempt_dir / "render_report.json"),
                "region_crop_plan": valid_region_crop_plan(),
                "crop_manifest": {"status": "empty-or-not-required", "reason": "no_source_crops_required"},
                "editable_text_count": 8,
                "native_shapes_count": 4,
                "source_crop_count": 3,
                "visual_qa": {"known_visual_risks": []},
                "failure_reason": None,
                "limitations": [],
                "recommended_next_round_change": "",
                "worker_identity": {
                    "worker_thread_id": "worker-1",
                    "execution_model": "delegated-real-per-slide-worker",
                    "slide_scope": [1],
                },
                "ocr_runtime_status": "passed-text-usable",
                "single_page_engine_used": "ppt-to-editable-v3-1-preview",
                "single_page_engine_scripts_used": ["package_reconstruction_deck.py"],
            },
        )

        result = deck_controller.collect_slide_result(
            {
                "slide_index": 1,
                "required_reading": manifest["required_reading"],
                "qa_summary_required_fields": manifest["qa_summary_required_fields"],
                "required_progress_statuses": deck_controller.REQUIRED_PROGRESS_STATUSES,
                "ocr_dependency_status_at_dispatch": "passed",
                "ocr_runtime_status_at_dispatch": "passed-text-usable",
            },
            run_dir,
        )

        self.assertEqual(result["status"], "failed-retryable")
        self.assertIn("text_background", result["failure_reason"])

    def test_success_is_rejected_when_non_text_anchor_is_missing(self) -> None:
        run_dir = fresh_dir("missing-non-text-anchor-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")
        manifest = deck_controller.make_slide_job_manifest(
            "run-1",
            {
                "slide_index": 1,
                "pictures": [
                    {
                        "extracted_path": str(source),
                        "image_size": {"width_px": 1600, "height_px": 900},
                    }
                ],
            },
            run_dir,
            ocr_runtime_status="passed-text-usable",
        )
        attempt_dir = run_dir / "slide_jobs" / "slide_01" / "attempt_01"
        rendered_png = attempt_dir / "render" / "slide_01.png"
        rendered_png.parent.mkdir(parents=True, exist_ok=True)
        rendered_png.write_bytes(b"png")
        (attempt_dir / "slide.pptx").write_bytes(b"placeholder")
        write_json(
            attempt_dir / "progress.json",
            {
                "history": [
                    {"status": "started"},
                    {"status": "ocr_done"},
                    {"status": "plan_done"},
                    {"status": "packaging_done"},
                    {"status": "qa_done"},
                ]
            },
        )
        write_json(
            attempt_dir / "render_report.json",
            {
                "status": "rendered",
                "rendered_png": str(rendered_png),
            },
        )
        review_fields = visual_fidelity_review_fields()
        review_fields["non_text_anchor_preservation_review"] = {
            "status": "failed",
            "missing_required_anchors": True,
            "issues": [
                {
                    "issue": "missing-arrowhead",
                    "anchor_id": "p3-upper-right-route-arrow",
                }
            ],
        }
        write_json(
            attempt_dir / "qa_summary.json",
            {
                "status": "passed",
                "mode": "mixed-reconstruction",
                "references_read": [{"id": item["id"]} for item in manifest["required_reading"]],
                "fallback_references_read": [],
                "font_policy_applied": "microsoft-yahei-all-runs",
                **review_fields,
                "effects_cleared": True,
                "rendered_png_reviewed": True,
                "render_report": str(attempt_dir / "render_report.json"),
                "region_crop_plan": valid_region_crop_plan(),
                "crop_manifest": {"status": "empty-or-not-required", "reason": "no_source_crops_required"},
                "editable_text_count": 8,
                "native_shapes_count": 4,
                "source_crop_count": 3,
                "visual_qa": {"known_visual_risks": []},
                "failure_reason": None,
                "limitations": [],
                "recommended_next_round_change": "",
                "worker_identity": {
                    "worker_thread_id": "worker-1",
                    "execution_model": "delegated-real-per-slide-worker",
                    "slide_scope": [1],
                },
                "ocr_runtime_status": "passed-text-usable",
                "single_page_engine_used": "ppt-to-editable-v3-1-preview",
                "single_page_engine_scripts_used": ["package_reconstruction_deck.py"],
            },
        )

        result = deck_controller.collect_slide_result(
            {
                "slide_index": 1,
                "required_reading": manifest["required_reading"],
                "qa_summary_required_fields": manifest["qa_summary_required_fields"],
                "required_progress_statuses": deck_controller.REQUIRED_PROGRESS_STATUSES,
                "ocr_dependency_status_at_dispatch": "passed",
                "ocr_runtime_status_at_dispatch": "passed-text-usable",
            },
            run_dir,
        )

        self.assertEqual(result["status"], "failed-retryable")
        self.assertIn("non_text_anchor", result["failure_reason"])


if __name__ == "__main__":
    unittest.main()
