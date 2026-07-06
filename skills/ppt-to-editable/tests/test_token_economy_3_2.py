import importlib.util
import json
import shutil
import tempfile
import unittest
from pathlib import Path


SCRIPT_PATH = (
    Path(__file__).resolve().parents[1]
    / "scripts"
    / "deck_controller.py"
)

spec = importlib.util.spec_from_file_location("deck_controller", SCRIPT_PATH)
deck_controller = importlib.util.module_from_spec(spec)
assert spec.loader is not None
spec.loader.exec_module(deck_controller)


TMP_ROOT = Path(tempfile.gettempdir()) / "ppt_to_editable_3_2_token_economy_tests"
TMP_ROOT.mkdir(parents=True, exist_ok=True)


def fresh_dir(name: str) -> Path:
    path = TMP_ROOT / name
    if path.exists():
        shutil.rmtree(path)
    path.mkdir(parents=True)
    return path


def write_json(path: Path, data: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def valid_gates_for(source_pptx: Path, conversion_mode: str, worker_mode: str) -> dict:
    return {
        "schema_version": deck_controller.GATES_SCHEMA_VERSION,
        "source_pptx_sha256": deck_controller.sha256_file(source_pptx),
        "conversion_mode_gate": {
            "user_choice": conversion_mode,
            "user_response_quote": "user selected token economy mode",
        },
        "ocr_runtime_gate": {
            "probe_ocr_runtime_status": "passed-text-usable",
            "setup_required": False,
            "auto_passed_existing_runtime": True,
            "explained_to_user": False,
            "user_confirmed_setup": False,
            "user_accepted_limited_run": False,
            "user_response_quote": "",
        },
        "scope_and_worker_gate": {
            "user_choice": "all",
            "slides": "",
            "worker_mode": worker_mode,
            "worker_execution_approved": True,
            "external_codex_worker_approved": False,
            "token_cost_acknowledged": True,
            "content_sharing_acknowledged": True,
            "user_response_quote": "convert all pages",
        },
    }


class TokenEconomyManifestTests(unittest.TestCase):
    def test_probe_reports_token_and_pro_mode_notes_without_internal_worker_terms(self) -> None:
        from PIL import Image
        from pptx import Presentation

        run_dir = fresh_dir("probe-token-pro-notes-run")
        source_pptx = run_dir / "source.pptx"
        image_path = run_dir / "source.png"
        Image.new("RGB", (1600, 900), "white").save(image_path)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        prs.save(source_pptx)

        original_get_ocr_dependency_status = deck_controller.get_ocr_dependency_status

        def fake_ocr_status(*_args, **_kwargs):
            return {
                "status": "passed",
                "ocr_runtime_status": "passed-text-usable",
                "checks": {},
                "runtime": {},
            }

        deck_controller.get_ocr_dependency_status = fake_ocr_status
        try:
            result = deck_controller.probe_deck(source_pptx)
        finally:
            deck_controller.get_ocr_dependency_status = original_get_ocr_dependency_status

        mode_gate = result["user_confirmation_required"]["conversion_mode_gate"]
        scope_gate = result["user_confirmation_required"]["scope_and_worker_gate"]
        self.assertIn("token", mode_gate["token_economy_note"])
        self.assertIn("token", scope_gate["token_cost_note"])
        self.assertNotIn("Pro", scope_gate["question"])
        self.assertNotIn("Codex Pro", scope_gate["token_cost_note"])
        self.assertNotIn("app-native", mode_gate["question"])
        self.assertNotIn("worker_mode", scope_gate["question"])

    def test_sequential_token_saving_mode_is_retired(self) -> None:
        run_dir = fresh_dir("sequential-gates-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        write_json(
            gates_file,
            valid_gates_for(
                source_pptx,
                conversion_mode="multi-page-sequential-token-saving",
                worker_mode="sequential",
            ),
        )

        with self.assertRaisesRegex(ValueError, "retired|not supported|下线"):
            deck_controller.load_and_validate_gates(
                gates_file,
                source_pptx,
                slides=None,
                all_slides=True,
                allow_ocr_not_ready=False,
            )

    def test_sequential_worker_mode_is_retired(self) -> None:
        run_dir = fresh_dir("sequential-rejects-app-native-run")
        source_pptx = run_dir / "source.pptx"
        source_pptx.write_bytes(b"placeholder")
        gates_file = run_dir / "gates.json"
        write_json(
            gates_file,
            valid_gates_for(
                source_pptx,
                conversion_mode="multi-agent-high-quality",
                worker_mode="sequential",
            ),
        )

        with self.assertRaisesRegex(ValueError, "sequential|retired|not supported|下线"):
            deck_controller.load_and_validate_gates(
                gates_file,
                source_pptx,
                slides=None,
                all_slides=True,
                allow_ocr_not_ready=False,
            )

    def test_sequential_slide_job_creation_is_retired(self) -> None:
        run_dir = fresh_dir("sequential-manifest-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")

        with self.assertRaisesRegex(ValueError, "sequential|retired|not supported|下线"):
            deck_controller.make_slide_job_manifest(
                "run-3-2",
                {
                    "slide_index": 1,
                    "pictures": [
                        {
                            "extracted_path": str(source),
                            "image_size": {"width_px": 1600, "height_px": 900},
                        }
                    ],
                },
                run_dir,
                ocr_runtime_status="passed-text-usable",
                worker_mode="sequential",
                conversion_mode="multi-page-sequential-token-saving",
            )

    def test_sequential_dispatch_report_is_ready_without_external_worker(self) -> None:
        dispatch_spec = importlib.util.spec_from_file_location(
            "prepare_worker_dispatch",
            Path(__file__).resolve().parents[1] / "scripts" / "prepare_worker_dispatch.py",
        )
        dispatch = importlib.util.module_from_spec(dispatch_spec)
        assert dispatch_spec.loader is not None
        dispatch_spec.loader.exec_module(dispatch)

        run_dir = fresh_dir("sequential-dispatch-run")
        write_json(
            run_dir / "deck_manifest.json",
            {
                "slide_count": 1,
                "selected_slide_indices": [1],
                "slides": [{"slide_index": 1, "job_manifest": str(run_dir / "slide_jobs" / "slide_01" / "slide_job_manifest.json")}],
                "user_confirmation_gates": {
                    "conversion_mode_gate": {"user_choice": "multi-agent-high-quality"},
                    "scope_and_worker_gate": {
                        "worker_mode": "sequential",
                        "worker_execution_approved": True,
                        "token_cost_acknowledged": True,
                        "content_sharing_acknowledged": True,
                    },
                },
            },
        )

        report = dispatch.build_dispatch_report(run_dir, worker_mode="sequential")
        self.assertEqual(report["status"], "blocked-invalid-worker-mode")
        self.assertTrue(report["issues"])

    def test_slide_job_uses_compact_protocol_as_default_required_reading(self) -> None:
        run_dir = fresh_dir("compact-protocol-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")

        manifest = deck_controller.make_slide_job_manifest(
            "run-3-2",
            {
                "slide_index": 1,
                "pictures": [
                    {
                        "extracted_path": str(source),
                        "image_size": {"width_px": 1600, "height_px": 900},
                    }
                ],
            },
            run_dir,
            ocr_runtime_status="passed-text-usable",
        )

        compact_protocol = run_dir / "WORKER_COMPACT_PROTOCOL.md"
        self.assertTrue(compact_protocol.exists())
        self.assertIn("v3 Preview", compact_protocol.read_text(encoding="utf-8"))

        must_read_ids = [
            item["id"] for item in manifest["required_reading"] if item.get("must_read")
        ]
        self.assertEqual(
            must_read_ids,
            [
                "worker_compact_protocol",
                "slide_job_manifest",
                "page_conversion_contract",
                "crop_manifest_contract",
            ],
        )

        fallback_ids = [
            item["id"]
            for item in manifest["fallback_reading"]
            if item.get("read_when_blocked_or_needed")
        ]
        self.assertIn("single_page_engine_skill_3_1", fallback_ids)
        self.assertIn("single_page_engine_reconstruction_qa", fallback_ids)
        self.assertIn("visual_fidelity_execution", fallback_ids)

        task_text = (run_dir / "slide_jobs" / "slide_01" / "AGENT_TASK.md").read_text(
            encoding="utf-8"
        )
        self.assertIn("WORKER_COMPACT_PROTOCOL.md", task_text)
        self.assertIn("Fallback Reading", task_text)
        self.assertNotIn("Read every file in **Required Reading**", task_text)

    def test_reference_loading_policy_records_fallback_reference_budget(self) -> None:
        run_dir = fresh_dir("reference-policy-run")
        source = run_dir / "source.png"
        source.write_bytes(b"placeholder")

        manifest = deck_controller.make_slide_job_manifest(
            "run-3-2",
            {
                "slide_index": 2,
                "pictures": [
                    {
                        "extracted_path": str(source),
                        "image_size": {"width_px": 1600, "height_px": 900},
                    }
                ],
            },
            run_dir,
            ocr_runtime_status="passed-text-usable",
        )

        policy = manifest["reference_loading_policy"]
        self.assertEqual(policy["default_mode"], "compact-protocol-first")
        self.assertEqual(
            policy["must_read_ids"],
            [
                "worker_compact_protocol",
                "slide_job_manifest",
                "page_conversion_contract",
                "crop_manifest_contract",
            ],
        )
        self.assertIn("single_page_engine_skill_3_1", policy["fallback_reading_ids"])
        self.assertIn("fallback_references_read", manifest["qa_summary_required_fields"])


if __name__ == "__main__":
    unittest.main()
